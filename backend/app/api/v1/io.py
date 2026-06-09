"""IO Router — modality-agnostic kernel endpoints (H5.3).

Endpoints:
- POST /chat/voice/transcribe  — Audio to text via Whisper
- POST /chat/voice/synthesize  — Text to speech via ElevenLabs
- POST /chat/documents/parse   — PDF/CSV/JSON parsing
- POST /chat/code/execute      — Sandboxed code execution
"""

from __future__ import annotations

import base64
import logging
import os
import time
from typing import TYPE_CHECKING

from fastapi import APIRouter, Depends, HTTPException, status

from app.api.deps import get_current_user
from app.models.io_models import (
    CodeExecuteRequest,
    CodeExecuteResponse,
    DocumentParseRequest,
    DocumentParseResponse,
    VoiceSynthesizeRequest,
    VoiceSynthesizeResponse,
    VoiceTranscribeRequest,
    VoiceTranscribeResponse,
)

if TYPE_CHECKING:
    from app.models.user import User

router = APIRouter(prefix="/chat", tags=["io"])
logger = logging.getLogger(__name__)


# ── Voice: Transcribe (Speech-to-Text via Whisper) ──────────────────


@router.post("/voice/transcribe", response_model=VoiceTranscribeResponse)
async def voice_transcribe(
    payload: VoiceTranscribeRequest,
    user: User = Depends(get_current_user),
):
    """Transcribe audio to text using Whisper (local or OpenAI API)."""
    logger.info("voice_transcribe: user=%s language=%s", user.id, payload.language)
    try:
        from app.tools.speech_to_text_transcriber import SpeechToTextTranscriberTool

        tool = SpeechToTextTranscriberTool()
        result = await tool.execute(
            {
                "data": payload.audio_data,
                "url": payload.audio_url,
                "language": payload.language,
                "response_format": "verbose_json",
            }
        )

        if result.status.value != "success":
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=result.error or "Transcription failed",
            )

        data = result.data or {}
        segments = data.get("segments", [])
        duration = data.get("duration", 0.0)

        return VoiceTranscribeResponse(
            text=data.get("text", ""),
            language=data.get("language"),
            duration_seconds=duration,
            segments=[
                {
                    "start": s.get("start", 0),
                    "end": s.get("end", 0),
                    "text": s.get("text", ""),
                    "confidence": s.get("confidence", 0),
                }
                for s in (segments or [])
            ],
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Voice transcription failed for user %s", user.id)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Transcription failed",
        )


# ── Voice: Synthesize (Text-to-Speech via ElevenLabs) ───────────────


@router.post("/voice/synthesize", response_model=VoiceSynthesizeResponse)
async def voice_synthesize(
    payload: VoiceSynthesizeRequest,
    user: User = Depends(get_current_user),
):
    """Convert text to speech using ElevenLabs TTS."""
    logger.info(
        "voice_synthesize: user=%s voice=%s text_len=%d",
        user.id,
        payload.voice_id,
        len(payload.text),
    )
    try:
        from app.tools.elevenlabs_tts import ElevenLabsTTSTool

        tool = ElevenLabsTTSTool()
        result = await tool.execute(
            {
                "text": payload.text,
                "voice_id": payload.voice_id,
                "model_id": "eleven_turbo_v2_5",
                "output_format": "mp3_44100_128",
                "save_to_storage": True,
            }
        )

        if result.status.value != "success":
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=result.error or "TTS synthesis failed",
            )

        data = result.data or {}
        audio_path = data.get("audio_path", "")
        audio_base64 = ""

        # Read the audio file and encode as base64 for the response
        if audio_path and os.path.exists(audio_path):
            with open(audio_path, "rb") as f:
                audio_bytes = f.read()
            audio_base64 = base64.b64encode(audio_bytes).decode("utf-8")

        return VoiceSynthesizeResponse(
            audio_base64=audio_base64 if not data.get("audio_url") else None,
            audio_url=data.get("audio_url"),
            format=data.get("format", "mp3"),
            duration_seconds=data.get("duration_seconds", 0.0),
            voice_id=data.get("voice_id", payload.voice_id),
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.exception("TTS synthesis failed for user %s", user.id)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Speech synthesis failed",
        )


# ── Documents: Parse (PDF, CSV, JSON) ───────────────────────────────


@router.post("/documents/parse", response_model=DocumentParseResponse)
async def document_parse(
    payload: DocumentParseRequest,
    user: User = Depends(get_current_user),
):
    """Parse a document (PDF, CSV, JSON) and return structured content."""
    logger.info(
        "document_parse: user=%s file=%s mime=%s",
        user.id,
        payload.filename,
        payload.mime_type,
    )
    try:
        file_data = payload.file_data
        filename = payload.filename
        mime_type = payload.mime_type

        # Resolve file data
        if not file_data and payload.file_url:
            import httpx

            async with httpx.AsyncClient(timeout=30) as client:
                resp = await client.get(payload.file_url)
                resp.raise_for_status()
                file_data = base64.b64encode(resp.content).decode("utf-8")
                if not filename:
                    filename = payload.file_url.rsplit("/", 1)[-1]

        if not file_data:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="No file data or URL provided",
            )

        raw_bytes = base64.b64decode(file_data)

        # H5.4: File size limit (50MB) to prevent abuse
        MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB
        if len(raw_bytes) > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File exceeds maximum size of 50MB (got {len(raw_bytes) / (1024 * 1024):.1f}MB)",
            )
        text_content = ""
        structured_data = None
        parse_error = None

        # Determine parser based on mime type or filename extension
        ext = os.path.splitext(filename)[1].lower() if filename else ""
        is_pdf = mime_type == "application/pdf" or ext == ".pdf"
        is_csv = mime_type == "text/csv" or ext == ".csv"
        is_json = mime_type == "application/json" or ext == ".json"
        is_pptx = (
            mime_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or ext == ".pptx"
        )
        is_docx = (
            mime_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or ext == ".docx"
        )

        if is_pdf:
            try:
                import io

                import PyPDF2

                reader = PyPDF2.PdfReader(io.BytesIO(raw_bytes))
                page_count = len(reader.pages)
                text_parts = []
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text_parts.append(t)
                text_content = "\n\n".join(text_parts)
                return DocumentParseResponse(
                    filename=filename,
                    mime_type=mime_type,
                    text_content=text_content[:100_000],
                    page_count=page_count,
                )
            except ImportError:
                parse_error = "PyPDF2 not installed"
            except Exception as e:
                parse_error = f"PDF parse error: {e}"

        if is_csv:
            try:
                import csv
                import io

                reader = csv.DictReader(
                    io.StringIO(raw_bytes.decode("utf-8", errors="replace"))
                )
                rows = list(reader)
                text_content = "\n".join(
                    ", ".join(f"{k}: {v}" for k, v in row.items()) for row in rows[:500]
                )
                # H5.4: Include actual row data (first 500 rows) + truncation flag
                structured_data = {
                    "columns": list(rows[0].keys()) if rows else [],
                    "row_count": len(rows),
                    "rows": rows[:500],
                    "truncated": len(rows) > 500,
                }
                return DocumentParseResponse(
                    filename=filename,
                    mime_type=mime_type,
                    text_content=text_content[:100_000],
                    structured_data=structured_data,
                )
            except Exception as e:
                parse_error = f"CSV parse error: {e}"

        elif is_json:
            try:
                import json

                data = json.loads(raw_bytes.decode("utf-8", errors="replace"))
                text_content = json.dumps(data, indent=2, default=str)[:100_000]
                structured_data = data
                return DocumentParseResponse(
                    filename=filename,
                    mime_type=mime_type,
                    text_content=text_content[:100_000],
                    structured_data=structured_data,
                )
            except Exception as e:
                parse_error = f"JSON parse error: {e}"

        if is_pptx:
            try:
                import io as _io

                from pptx import Presentation

                prs = Presentation(_io.BytesIO(raw_bytes))
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_parts = [f"--- Slide {i} ---"]
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_parts.append(shape.text.strip())
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                row_text = " | ".join(
                                    cell.text.strip() for cell in row.cells
                                )
                                slide_parts.append(row_text)
                    slides_text.append("\n".join(slide_parts))
                text_content = "\n\n".join(slides_text)
                return DocumentParseResponse(
                    filename=filename,
                    mime_type=mime_type,
                    text_content=text_content[:100_000],
                    page_count=len(prs.slides),
                )
            except ImportError:
                parse_error = "python-pptx not installed"
            except Exception as e:
                parse_error = f"PPTX parse error: {e}"

        if is_docx:
            try:
                import io as _io

                from docx import Document

                doc = Document(_io.BytesIO(raw_bytes))
                paragraphs = []
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if text:
                        # Add heading markers
                        if para.style.name.startswith("Heading"):
                            level = para.style.name.replace("Heading ", "")
                            paragraphs.append(f"{'#' * int(level)} {text}")
                        else:
                            paragraphs.append(text)
                # Also extract tables
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            paragraphs.append(row_text)
                text_content = "\n\n".join(paragraphs)
                return DocumentParseResponse(
                    filename=filename,
                    mime_type=mime_type,
                    text_content=text_content[:100_000],
                    page_count=len(doc.sections),
                )
            except ImportError:
                parse_error = "python-docx not installed"
            except Exception as e:
                parse_error = f"DOCX parse error: {e}"

        elif not parse_error:
            parse_error = f"Unsupported file type: {mime_type or ext}"

        # Fallback: return raw text
        try:
            text_content = raw_bytes.decode("utf-8", errors="replace")[:100_000]
        except Exception:
            text_content = f"[Binary file: {len(raw_bytes)} bytes, {mime_type}]"

        return DocumentParseResponse(
            filename=filename,
            mime_type=mime_type,
            text_content=text_content[:100_000],
            error=parse_error,
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Document parsing failed for user %s", user.id)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Document parsing failed",
        )


# ── Code: Execute ───────────────────────────────────────────────────


@router.post("/code/execute", response_model=CodeExecuteResponse)
async def code_execute(
    payload: CodeExecuteRequest,
    user: User = Depends(get_current_user),
):
    """Execute a code cell using the production sandbox tools (Python or Node.js).

    Uses PythonSandboxTool / NodeJsSandboxTool for proper resource limits,
    import denylists, and timeout handling (H5.4 production wiring).
    """
    logger.info(
        "code_execute: user=%s language=%s timeout=%ds",
        user.id,
        payload.language,
        payload.timeout_seconds,
    )
    start = time.monotonic()

    lang = (payload.language or "python").lower()

    # Route to appropriate sandbox tool (H5.4 production wiring)
    try:
        if lang in ("javascript", "js", "typescript", "ts", "node", "nodejs"):
            from app.tools.nodejs_sandbox import NodeJsSandboxTool

            tool = NodeJsSandboxTool()
            result = await tool.execute(
                {
                    "code": payload.code,
                    "timeout_seconds": payload.timeout_seconds,
                }
            )
        else:
            from app.tools.python_sandbox import PythonSandboxTool

            tool = PythonSandboxTool()
            result = await tool.execute(
                {
                    "code": payload.code,
                    "timeout_seconds": payload.timeout_seconds,
                }
            )
    except ImportError:
        raise HTTPException(
            status_code=status.HTTP_501_NOT_IMPLEMENTED,
            detail=f"{lang} sandbox is not available",
        )
    except Exception as e:
        elapsed_ms = (time.monotonic() - start) * 1000
        logger.exception("Sandbox execution failed for user %s", user.id)
        return CodeExecuteResponse(
            success=False,
            execution_time_ms=round(elapsed_ms, 1),
            error="Code execution failed",
        )

    elapsed_ms = (time.monotonic() - start) * 1000

    if result.status.value != "success":
        return CodeExecuteResponse(
            success=False,
            execution_time_ms=round(elapsed_ms, 1),
            error=result.error or "Code execution failed",
        )

    data = result.data or {}
    stdout_val = str(data.get("stdout", data.get("output", "")))[:500_000]
    stderr_val = str(data.get("stderr", ""))[:100_000]
    return_code = data.get("return_code", data.get("returncode", 0))

    return CodeExecuteResponse(
        success=True,
        stdout=stdout_val,
        stderr=stderr_val,
        return_code=return_code,
        execution_time_ms=round(elapsed_ms, 1),
        error=stderr_val if return_code != 0 else None,
    )
