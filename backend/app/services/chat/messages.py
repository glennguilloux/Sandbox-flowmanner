# ─────────────────────────────────────────────────────────────────────────
# Auto-decomposed from app/services/chat_service.py (CARD 3 refactor).
# Part of the `chat` package. Sibling cross-references and original imports
# are preserved so behavior/signatures stay byte-for-byte identical.
# ─────────────────────────────────────────────────────────────────────────
from __future__ import annotations

import base64
import os
from datetime import UTC
from pathlib import Path
from typing import TYPE_CHECKING

from sqlalchemy import func, select

from app.models.chat import ChatBranch, ChatFile, ChatMessage, ChatThread
from app.models.phase4_models import UserFile
from app.services.chat_context import _inject_memory_context, _prune_messages_to_budget

from .prompts import _inject_web_search

if TYPE_CHECKING:
    from sqlalchemy.ext.asyncio import AsyncSession

__all__ = [
    "_get_model_preference",
    "_prepare_step_inject",
    "_process_attachments",
    "create_chat_file",
    "create_chat_message",
    "create_chat_message_fresh_session",
    "delete_chat_message",
    "get_chat_files",
    "get_chat_messages",
    "update_chat_message",
]


async def create_chat_message(
    db: AsyncSession,
    thread_id: int,
    role: str,
    content: str,
    *,
    user_id: int | None = None,
) -> ChatMessage:
    msg = ChatMessage(thread_id=thread_id, role=role, content=content)
    if user_id is not None:
        msg.user_id = user_id
    db.add(msg)
    await db.flush()
    await db.refresh(msg)
    return msg


async def create_chat_message_fresh_session(
    thread_id: int,
    role: str,
    content: str,
    *,
    user_id: int | None = None,
) -> ChatMessage:
    """Create a chat message using a fresh DB session.

    Used as a fallback when the caller's session has a dead connection
    (e.g. after long LLM streaming where idle-in-transaction kills the
    underlying asyncpg connection).  The session is committed by the
    context manager on successful exit.
    """
    from app.database import fresh_session

    async with fresh_session() as fresh_db:
        msg = ChatMessage(thread_id=thread_id, role=role, content=content)
        if user_id is not None:
            msg.user_id = user_id
        fresh_db.add(msg)
        await fresh_db.flush()
        await fresh_db.refresh(msg)
        return msg


async def update_chat_message(
    db: AsyncSession,
    message_id: int,
    content: str,
) -> ChatMessage | None:
    """Update a chat message's content and set edited_at timestamp."""
    from datetime import datetime

    result = await db.execute(select(ChatMessage).where(ChatMessage.id == message_id))
    message = result.scalar_one_or_none()
    if message is None:
        return None
    message.content = content
    message.edited_at = datetime.now(UTC)
    await db.flush()
    await db.refresh(message)
    return message


async def delete_chat_message(db: AsyncSession, message_id: int) -> bool:
    """Delete a chat message by ID."""
    result = await db.execute(select(ChatMessage).where(ChatMessage.id == message_id))
    message = result.scalar_one_or_none()
    if message is None:
        return False
    await db.delete(message)
    await db.flush()
    return True


async def get_chat_messages(
    db: AsyncSession,
    thread_id: int,
    *,
    offset: int = 0,
    limit: int = 50,
) -> tuple[list[ChatMessage], int]:
    count_q = select(func.count()).select_from(ChatMessage).where(ChatMessage.thread_id == thread_id)
    total = (await db.execute(count_q)).scalar() or 0
    items_q = (
        select(ChatMessage)
        .where(ChatMessage.thread_id == thread_id)
        .order_by(ChatMessage.created_at.asc())
        .offset(offset)
        .limit(limit)
    )
    items = list((await db.execute(items_q)).scalars().all())
    return items, total


async def get_chat_files(
    db: AsyncSession,
    thread_id: int,
) -> list[ChatFile]:
    result = await db.execute(select(ChatFile).where(ChatFile.chat_id == thread_id))
    return list(result.scalars().all())


async def create_chat_file(
    db: AsyncSession,
    thread_id: int,
    filename: str,
    mime_type: str,
    path: str,
    size_bytes: int,
) -> ChatFile:
    file = ChatFile(
        chat_id=thread_id,
        filename=filename,
        mime_type=mime_type,
        path=path,
        size_bytes=size_bytes,
    )
    db.add(file)
    await db.flush()
    await db.refresh(file)
    return file


def _get_model_preference(thread: ChatThread) -> str | None:
    if thread.metadata_ and isinstance(thread.metadata_, dict):
        return thread.metadata_.get("model_preference")
    return None


async def _process_attachments(
    db: AsyncSession, messages: list[dict], attachments: list[dict], model: str
) -> list[dict]:
    is_vision_model = not model.startswith("llamacpp/")

    for att in attachments:
        file_id = att.get("file_id", "")
        att_type = att.get("type", "file")
        filename = att.get("filename", "unknown")

        result = await db.execute(select(UserFile).where(UserFile.id == file_id))
        db_file = result.scalar_one_or_none()
        if not db_file or not db_file.storage_path or not os.path.exists(db_file.storage_path):
            continue

        if att_type == "image" and is_vision_model:
            raw_bytes = Path(db_file.storage_path).read_bytes()
            b64 = base64.b64encode(raw_bytes).decode("utf-8")
            content_type = db_file.content_type or "image/png"

            if messages and messages[-1].get("role") == "user":
                existing = messages[-1].get("content", "")
                if isinstance(existing, str):
                    messages[-1]["content"] = [
                        {"type": "text", "text": existing},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{content_type};base64,{b64}"},
                        },
                    ]
                else:
                    messages[-1]["content"].append(
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{content_type};base64,{b64}"},
                        }
                    )

        elif att_type == "image" and not is_vision_model:
            if messages and messages[-1].get("role") == "user":
                existing = messages[-1].get("content", "")
                if isinstance(existing, str):
                    messages[-1]["content"] = f"{existing}\n\n[Image: {filename}]"

        elif att_type == "file":
            try:
                content_type = (db_file.content_type or "").lower()
                filename_lower = (db_file.filename or "").lower()

                if "pdf" in content_type:
                    import pdfplumber

                    with pdfplumber.open(db_file.storage_path) as pdf:
                        pages_text = []
                        for page in pdf.pages:
                            t = page.extract_text() or ""
                            pages_text.append(t)
                        file_text = "\n\n".join(pages_text)

                elif "wordprocessingml" in content_type or ".docx" in filename_lower or "msword" in content_type:
                    from docx import Document

                    doc = Document(db_file.storage_path)
                    file_text = "\n".join(p.text for p in doc.paragraphs)

                elif "spreadsheetml" in content_type or ".xlsx" in filename_lower or "excel" in content_type:
                    import openpyxl

                    wb = openpyxl.load_workbook(db_file.storage_path, read_only=True, data_only=True)
                    rows = []
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        sheet_rows = []
                        for row in ws.iter_rows(values_only=True):
                            line = "\t".join(str(c) if c is not None else "" for c in row)
                            sheet_rows.append(line)
                        rows.append(f"=== Sheet: {sheet} ===\n" + "\n".join(sheet_rows))
                    file_text = "\n\n".join(rows)
                    wb.close()

                elif "presentationml" in content_type or ".pptx" in filename_lower or "powerpoint" in content_type:
                    from pptx import Presentation

                    prs = Presentation(db_file.storage_path)
                    slides_text = []
                    for i, slide in enumerate(prs.slides, 1):
                        slide_texts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                slide_texts.append(shape.text)
                            elif shape.has_table:
                                table = shape.table
                                slide_texts.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
                        slides_text.append(f"=== Slide {i} ===\n" + "\n".join(slide_texts))
                    file_text = "\n\n".join(slides_text)

                else:
                    file_text = Path(db_file.storage_path).read_text(encoding="utf-8", errors="replace")
            except Exception:
                continue
            limit = 10000
            truncated = file_text[:limit]
            if len(file_text) > limit:
                truncated += "\n... (truncated)"
            context_msg = f"[Attached file: {filename}]\n{truncated}\n[End of attached file]"
            messages.insert(-1, {"role": "user", "content": context_msg})

    return messages


async def _prepare_step_inject(
    messages: list[dict],
    *,
    memory_claims: list | None = None,
    web_search: bool = False,
    content: str | None = None,
    steps: list | None = None,
    should_inject: bool | None = None,
) -> tuple[list[dict], list[dict]]:
    """Ordered context-injection closure for chat (SPIKE — ADR-002).

    Mirrors trigger.dev ``chat.agent``'s single ``prepareStep`` primitive:
    all context sources are merged into ``messages`` in a fixed order
    (memory -> web search) and the function returns both the mutated message
    list and a list of receipt events.

    Today this is called once per turn, pre-LLM, with ``steps=None``
    (single-shot chat). The future re-entrant turn loop will call it at each
    AI-SDK step boundary with a non-empty ``steps`` list and a ``should_inject``
    gate, matching trigger.dev's step-boundary injection model. The ordering is
    enforced here by code structure, not convention.

    Behavior when the flag is off is identical to the legacy inline path
    (see ``_stream_message_to_llm_body``); this helper is only reached when
    ``settings.CHAT_PREPARE_STEP_HOOK_ENABLED`` is True.

    Returns:
        (messages_with_context, injected_events) where each event is an SSE
        payload dict the caller yields so the frontend can reconcile
        injected-vs-queued context.
    """
    injected_events: list[dict] = []
    result = list(messages)

    # 1) Memory context (ordered first, like prepareStep's compaction->steering)
    if memory_claims:
        result = _inject_memory_context(result, memory_claims)
        injected_events.append(
            {
                "type": "injected",
                "source": "memory",
                "count": len(memory_claims),
            }
        )

    # 2) Web search context
    if web_search and content is not None:
        result = await _inject_web_search(result, content)
        injected_events.append(
            {
                "type": "injected",
                "source": "web_search",
                "query": content,
            }
        )

    # (Future step-boundary hook point) if steps is not None, apply a
    # should_inject gate and additional per-step sources here.
    return result, injected_events
