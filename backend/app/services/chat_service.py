from __future__ import annotations

import asyncio
import base64
import contextlib
import json
import logging
import os
import time
from collections import deque
from collections.abc import AsyncGenerator
from datetime import UTC
from pathlib import Path
from typing import TYPE_CHECKING, cast

if TYPE_CHECKING:
    from collections.abc import AsyncGenerator

from openai import AsyncOpenAI
from sqlalchemy import func, select

from app.models.chat import ChatBranch, ChatFile, ChatMessage, ChatThread
from app.models.phase4_models import UserFile

logger = logging.getLogger(__name__)

from app.config import settings
from app.services.chat_context import _inject_memory_context, _prune_messages_to_budget
from app.services.llm_providers import (
    _LLM_API_BASE,
    _LLM_API_KEY,
    _LLM_MODEL,
    OPENAI_PROVIDER_FAMILIES,
    PROVIDER_MAP,
    _detect_provider_from_key,
    _get_base_url_for_provider,
    _get_provider_for_model,
    _get_upstream_model_name,
    _normalize_provider,
    _providers_compatible,
    _resolve_provider,
)
from app.services.memory_citation_service import (
    build_citation_event,
    build_recall_used_event,
)
from app.services.personal_memory_extractor import (
    PersonalMemoryExtractor,
    RegexPersonalMemoryExtractor,
)

if TYPE_CHECKING:
    from sqlalchemy.ext.asyncio import AsyncSession

    from app.models.personal_memory_models import PersonalMemoryClaim


async def _safe_fire_and_forget(coro, *, label: str) -> None:
    """Wrap a fire-and-forget coroutine so exceptions are logged, not silently dropped."""
    try:
        await coro
    except Exception:
        logger.exception("fire-and-forget task failed: %s", label)


async def _lookup_stored_byok_key(
    db: AsyncSession,
    user_id: int,
    provider_hint: str | None = None,
) -> tuple[str | None, str | None]:
    """Look up the best stored BYOK key for a user.

    Returns (api_key, base_url) or (None, None) if no key found.
    If provider_hint is given, prefers a key whose provider matches.
    Falls back to any active key if no provider match.
    """
    logger.debug("BYOK lookup: user=%s provider=%s", user_id, provider_hint)
    try:
        from sqlalchemy import select

        from app.models.byok_models import UserAPIKey

        stmt = select(UserAPIKey).where(UserAPIKey.user_id == user_id).where(UserAPIKey.is_active == True)
        result = await db.execute(stmt)
        keys = list(result.scalars().all())
        if not keys:
            return None, None

        # Prefer a key whose provider matches the hint
        if provider_hint:
            normalized_hint = _normalize_provider(provider_hint)
            for k in keys:
                if _normalize_provider(k.provider) == normalized_hint:
                    logger.info(
                        "BYOK lookup FOUND: user=%s provider=%s key_id=%s",
                        user_id,
                        provider_hint,
                        k.id,
                    )
                    return k.get_api_key(), k.base_url

            # No provider match found — do NOT fall back to a random key.
            # Returning a key for the wrong provider causes 400 errors
            # (e.g. sending demo-llm to OpenRouter instead of glennguilloux).
            logger.info(
                "BYOK lookup NO MATCH: user=%s provider=%s available=%s",
                user_id,
                provider_hint,
                [k.provider for k in keys],
            )
            return None, None

        # Fall back to the first active key
        logger.debug(
            "BYOK lookup FALLBACK: user=%s provider=%s -> using first key key_id=%s",
            user_id,
            provider_hint,
            keys[0].id,
        )
        return keys[0].get_api_key(), keys[0].base_url
    except Exception:
        logger.warning(
            "BYOK lookup ERROR: user=%s provider=%s",
            user_id,
            provider_hint,
            exc_info=True,
        )
        return None, None


def _validate_byok_key_matches_model(user_api_key: str | None, model_id: str) -> str | None:
    """Validate that BYOK key matches the requested model provider.

    Returns None if valid.
    Returns error message if mismatch:
    - For llamacpp/*: always valid (ignore key)
    - For openai_compatible/*: accept any OpenAI-family key or unknown keys
    - For other providers: check if key matches model provider
    """
    if not user_api_key or not model_id:
        return None

    model_provider = _get_provider_for_model(model_id)

    if model_provider == "llamacpp":
        return None

    key_provider = _detect_provider_from_key(user_api_key)

    if not _providers_compatible(key_provider, model_provider):
        return f"Provider mismatch: model '{model_id}' requires {model_provider.title()}, but X-User-API-Key appears to be for {key_provider.title()}"

    return None


# Sandboxd preview guidance appended to system prompt when enabled
_SANDBOXD_SYSTEM_GUIDANCE = """

## Live Preview Tools (sandboxd)

When the user asks you to build something visual (landing page, dashboard,
chart, tool, app, or any HTML/CSS/JS project), use the sandboxd tools to
create a live preview. Follow this workflow **exactly**:

1. **sandboxd_preview** — call with `{}` (no arguments) to create a new sandbox.
   Save the returned `sandbox_id` for ALL subsequent calls.
   ⚠️  CRITICAL: sandboxd_preview returns sandbox metadata ONLY (id, status).
   It does NOT return a usable app preview URL.  The sandbox runtime URL (port 3000)
   is empty — do NOT show it to the user.  The preview URL comes from sandboxd_serve.
2. **sandboxd_file_write** — write your files. Pass `sandbox_id` and `path` and `content`.
   Example: `{"sandbox_id": "...", "path": "index.html", "content": "<!DOCTYPE html>..."}`
   Subdirectories are created automatically (e.g. `"css/style.css"` works fine).
   Write ALL files before calling sandboxd_serve.
3. **sandboxd_serve** — start a dev server and get the preview URL:
   `{"sandbox_id": "..."}`
   This starts a server on port 8081 (the default python-img template does not use 8080; some legacy templates like react-standard may)
   that serves from /home/sandbox/ where your files were written.
   This is the ONLY tool that returns the app preview URL.  ALWAYS call it
   after writing files, and ALWAYS present its returned URL to the user.
4. **sandboxd_file_read** — read a file back: `{"sandbox_id": "...", "path": "index.html"}`
5. **sandboxd_file_list** — list workspace files: `{"sandbox_id": "...", "path": ""}`

If you need to run custom commands (npm install, python scripts, etc.), use
**sandboxd_exec** with the `command` field (argv array).

The preview URL format is:
https://s-<sandbox_id>-8081.preview.flowmanner.com

The URL is publicly accessible. The sandbox stays alive for 35 minutes.

If any tool returns an error containing "container is not running" or
"not running", do NOT retry more than once — explain the error to the user."""

# Maximum number of tool-call rounds before forcing a text response
_MAX_TOOL_ROUNDS = settings.CHAT_MAX_TOOL_ROUNDS


from app.services.background_task_manager import background_task_manager
from app.services.sse_protocol import _CANVAS_UPDATE_TOOLS, _build_canvas_update

# -- SSE keepalive ping (Task 2b.2) -----------------------------------------
# Defeats proxy idle timeouts during long tool rounds. The keepalive must be
# timer-driven (every 15s) and independent of token cadence, because a single
# tool round can block the generator for >15s with no yield — a yield-gated
# ping would never fire and nginx would drop the idle connection.
_SSE_KEEPALIVE_INTERVAL = 15  # seconds

# SSE comment line. `: ` prefix makes it a comment the browser EventSource ignores.
_SSE_KEEPALIVE_PING = ": ping\n\n"


async def _sse_keepalive_timer(queue: asyncio.Queue, stop_event: asyncio.Event) -> None:
    """Emit `: ping` comments every ``_SSE_KEEPALIVE_INTERVAL`` seconds.

    Runs as a ``BackgroundTaskManager`` task. Pumps ``_SSE_KEEPALIVE_PING``
    into ``queue`` on a fixed cadence until ``stop_event`` is set or cancelled.
    Pure timer — never depends on LLM token cadence.
    """
    try:
        while not stop_event.is_set():
            try:
                await asyncio.wait_for(stop_event.wait(), timeout=_SSE_KEEPALIVE_INTERVAL)
            except TimeoutError:
                # Interval elapsed without cancellation — emit a keepalive.
                await queue.put(_SSE_KEEPALIVE_PING)
            else:
                # stop_event was set during the wait — exit cleanly.
                break
    except asyncio.CancelledError:
        # Expected on stream teardown; swallow so the manager logs nothing.
        raise


async def _sse_keepalive_merge(
    body: AsyncGenerator[str, None],
    queue: asyncio.Queue,
    stop_event: asyncio.Event,
) -> AsyncGenerator[str, None]:
    """Merge the LLM ``body`` stream with the timer-driven keepalive queue.

    Yields events from ``body`` as they arrive (record = "real activity" →
    reset the keepalive clock), and interleaves ``_SSE_KEEPALIVE_PING`` entries
    from ``queue`` that arrive while the body is idle. The keepalive timer is
    live for the whole stream, so a multi-second tool round still gets pings.
    """
    body_iter = body.__aiter__()
    # One in-flight anext coroutine at a time. Creating a fresh __anext__()
    # while the previous is pending raises "async generator is already running",
    # so we hold the future and only advance after consuming it.
    body_task = asyncio.ensure_future(body_iter.__anext__())
    while True:
        # Wait for the next body event, but wake at least every interval so
        # idle gaps still surface timer pings.
        done, _pending = await asyncio.wait({body_task}, timeout=_SSE_KEEPALIVE_INTERVAL)
        if body_task in done:
            try:
                event = body_task.result()
            except StopAsyncIteration:
                # Body is exhausted. Any ping the timer queued during the final
                # idle gap must NOT be appended after the real last body event,
                # so we break here without draining the queue. The timer is
                # cancelled on generator close.
                break
            # Advance to the next body event for the next iteration.
            body_task = asyncio.ensure_future(body_iter.__anext__())
            yield event
            continue

        # No body event within the interval → this is a genuine idle gap.
        # Surface exactly one queued ping (the timer self-rate-limits to one
        # per interval, but the queue may hold a straggler from a previous gap).
        # We only drain the queue in this branch, so pings never leak after the
        # real last body event.
        try:
            ping = queue.get_nowait()
        except asyncio.QueueEmpty:
            # Timer hasn't fired yet (scheduler jitter); emit the ping directly
            # so the idle gap is still covered even if the queue was drained.
            ping = _SSE_KEEPALIVE_PING
        yield ping


async def stream_message_to_llm(
    db: AsyncSession,
    thread_id: int,
    content: str,
    user_id: int,
    model_preference: str | None = None,
    user_api_key: str | None = None,
    user_base_url: str | None = None,
    model_id: str | None = None,
    attachments: list | None = None,
    web_search: bool | None = None,
) -> AsyncGenerator[str, None]:
    """Send a message to the LLM and stream the response via SSE.

    The returned generator drives its own timer-driven keepalive (every 15s,
    independent of token cadence) by spawning a ``BackgroundTaskManager`` task
    and merging its pings with the LLM event stream. The keepalive task is
    cancelled and awaited on generator close so no orphan task leaks.

    If user_api_key is provided, a per-request AsyncOpenAI client is created
    using that key (BYOK path). The key is used only for this call and discarded.
    If model_id is provided it overrides model_preference and the env default.
    If user_base_url is provided, it overrides the resolved base_url for custom providers.
    For llamacpp/* models, any BYOK key is ignored (llama.cpp doesn't use API keys).
    """
    _keepalive_queue: asyncio.Queue = asyncio.Queue()
    _keepalive_stop = asyncio.Event()
    _keepalive_task = background_task_manager.spawn(
        _sse_keepalive_timer(_keepalive_queue, _keepalive_stop),
        label="sse_keepalive",
    )
    try:
        async for event in _sse_keepalive_merge(
            cast(
                "AsyncGenerator[str, None]",
                _stream_message_to_llm_body(
                    db,
                    thread_id,
                    content,
                    user_id,
                    model_preference,
                    user_api_key,
                    user_base_url,
                    model_id,
                    attachments,
                    web_search,
                ),
            ),
            _keepalive_queue,
            _keepalive_stop,
        ):
            yield event
    finally:
        _keepalive_stop.set()
        _keepalive_task.cancel()
        with contextlib.suppress(asyncio.CancelledError, Exception):
            await _keepalive_task


async def _stream_message_to_llm_body(
    db: AsyncSession,
    thread_id: int,
    content: str,
    user_id: int,
    model_preference: str | None = None,
    user_api_key: str | None = None,
    user_base_url: str | None = None,
    model_id: str | None = None,
    attachments: list | None = None,
    web_search: bool | None = None,
):
    """Send a message to the LLM and stream the response via SSE.

    If user_api_key is provided, a per-request AsyncOpenAI client is created
    using that key (BYOK path). The key is used only for this call and discarded.
    If model_id is provided it overrides model_preference and the env default.
    If user_base_url is provided, it overrides the resolved base_url for custom providers.
    For llamacpp/* models, any BYOK key is ignored (llama.cpp doesn't use API keys).
    """
    # --- Pre-LLM setup: resolve provider, BYOK lookup, build messages ---
    collected_chunks = []
    _last_yield = time.monotonic()  # Task 1.2a: keepalive tracking
    raw_model = model_id or model_preference or _LLM_MODEL
    base_url, api_key, model = _resolve_provider(raw_model)

    mismatch_error = _validate_byok_key_matches_model(user_api_key, raw_model)
    if mismatch_error:
        now = time.monotonic()
        ping = _sse_keepalive(_last_yield, now)
        if ping:
            yield ping
        _last_yield = now
        yield json.dumps({"type": "error", "error": mismatch_error})
        return

    # --- Stored-key fallback: if no per-request key, look up user's stored key ---
    effective_user_key = user_api_key
    effective_base_url = user_base_url
    if not effective_user_key and db is not None:
        model_provider = _get_provider_for_model(raw_model)
        stored_key, stored_base = await _lookup_stored_byok_key(db, user_id, provider_hint=model_provider)
        if stored_key:
            effective_user_key = stored_key
            effective_base_url = stored_base

    if raw_model and raw_model.startswith("llamacpp/"):
        effective_user_key = None

    if effective_user_key:
        effective_base = effective_base_url or base_url
        client = AsyncOpenAI(api_key=effective_user_key, base_url=effective_base)
    elif base_url != _LLM_API_BASE or api_key != _LLM_API_KEY:
        client = AsyncOpenAI(api_key=api_key, base_url=base_url)
    else:
        client = _client

    # Circuit breaker protection
    from app.core.circuit_breaker import CircuitOpenError, get_circuit_breaker
    from app.core.metrics import record_llm_request

    provider_name = _get_provider_for_model(raw_model) or "deepseek"
    breaker = get_circuit_breaker(provider_name)
    llm_start = time.time()

    # Build messages and process all DB-dependent setup BEFORE saving user message.
    messages_for_llm = await _build_chat_messages(db, thread_id)
    # Append the new user message — it isn't in DB history yet
    messages_for_llm.append({"role": "user", "content": content})

    # ── Phase 5: resolve workspace_id for allowlist filtering ────────
    _thread_obj = await get_chat_thread(db, thread_id)
    _workspace_id = _thread_obj.workspace_id if _thread_obj else None

    # Phase 5: async tool list with workspace allowlist
    openai_tools = await _get_chat_openai_tools(db=db, workspace_id=_workspace_id)

    # ── Phase 2: pre-compute user scopes for tool authorization ──────
    # Resolve once before the tool loop so _execute_tool_call doesn't
    # need a DB lookup per tool invocation.
    _cached_user_scopes: set[str] | None = None
    _cached_user_role: str | None = None
    try:
        from app.models.user import User as UserModel

        _urow = (await db.execute(select(UserModel).where(UserModel.id == user_id))).scalar_one_or_none()
        if _urow is not None:
            _cached_user_scopes = set(getattr(_urow, "scopes", []) or [])
            _cached_user_role = getattr(_urow, "role", None)
    except Exception:
        logger.debug("stream_message_to_llm: scope pre-fetch failed (non-fatal)", exc_info=True)

    # ── T33 Stage 1: pre-LLM memory recall + injection ──────
    # Epic 2.2: freeze one recall_for_chat result per session (thread_id)
    # via the snapshot seam, instead of re-calling recall per message.
    # Cache miss -> exactly ONE recall_for_chat (seed query "") + store.
    # Cache hit -> reuse the frozen claims. Write-invalidation / TTL are
    # handled inside the seam. Injection shape below is unchanged.
    memory_recall_claims: list[PersonalMemoryClaim] = []
    if settings.CHAT_MEMORY_CITATIONS_ENABLED:
        thread_for_recall = await get_chat_thread(db, thread_id)
        if thread_for_recall and thread_for_recall.workspace_id:
            try:
                from app.services.memory_snapshot_service import (
                    get_or_capture_snapshot,
                )

                memory_recall_claims = await get_or_capture_snapshot(
                    db,
                    thread_id=thread_id,
                    user_id=user_id,
                    workspace_id=thread_for_recall.workspace_id,
                )
                # Q2-A/Q2-C: resolve the token-bounded, E23-B-ranked set.
                # Tier-0 constraints are protected; Tier-1 competitive claims
                # are dropped lowest-rank-first on overflow. Deterministic
                # (pure function) so the frozen snapshot stays reproducible.
                if memory_recall_claims:
                    from app.services.personal_memory_service import (
                        rank_and_budget_claims,
                    )

                    resolved, _dropped = rank_and_budget_claims(
                        memory_recall_claims,
                        token_budget=settings.CHAT_MEMORY_INJECTION_TOKEN_BUDGET,
                    )
                    if _dropped:
                        logger.info(
                            "stream_message_to_llm: memory budget dropped %d claim(s) " "(kept %d) for thread %s",
                            len(_dropped),
                            len(resolved),
                            thread_id,
                        )
                    memory_recall_claims = resolved
                    if memory_recall_claims:
                        logger.info(
                            "stream_message_to_llm: froze %d memory claims for thread %s",
                            len(memory_recall_claims),
                            thread_id,
                        )
            except Exception as recall_err:
                logger.warning(
                    "stream_message_to_llm: memory recall failed for thread %s, continuing without context: %s",
                    thread_id,
                    recall_err,
                )
                memory_recall_claims = []
        else:
            logger.info(
                "stream_message_to_llm: skipping memory recall for thread %s (workspace_id is null)",
                thread_id,
            )

    if attachments:
        messages_for_llm = await _process_attachments(db, messages_for_llm, attachments, raw_model)

    _prepare_step_injected_events: list[dict] = []
    if getattr(settings, "CHAT_PREPARE_STEP_HOOK_ENABLED", False):
        # Route context injection through the ordered prepareStep closure
        # (mirrors trigger.dev chat.agent prepareStep). Today this runs once
        # pre-LLM (single-shot chat, steps=None). The future re-entrant turn loop
        # will call it at each step boundary with a non-empty steps list +
        # shouldInject gate, matching trigger.dev's step-boundary injection model.
        messages_for_llm, _prepare_step_injected_events = await _prepare_step_inject(
            messages_for_llm,
            memory_claims=memory_recall_claims,
            web_search=bool(web_search),
            content=content,
        )
    else:
        # Legacy inline injection — unchanged behavior when the spike flag is off.
        if memory_recall_claims:
            messages_for_llm = _inject_memory_context(messages_for_llm, memory_recall_claims)
        if web_search:
            messages_for_llm = await _inject_web_search(messages_for_llm, content)

    # Save user message, commit, and immediately close the session to prevent
    # idle-in-transaction timeout during the LLM call
    await create_chat_message(db, thread_id, "user", content, user_id=user_id)
    await db.commit()
    await db.close()

    # Emit injection-receipt events so the frontend can reconcile
    # injected-vs-queued context (only when the prepareStep hook is enabled).
    for _ev in _prepare_step_injected_events:
        yield json.dumps(_ev)

    try:
        async with breaker.protect():
            full_response = ""
            accumulated_prompt_tokens = 0
            accumulated_completion_tokens = 0

            # ── Tool-calling loop ───────────────────────────────────
            for _round in range(_MAX_TOOL_ROUNDS):
                create_kwargs: dict = {
                    "model": model,
                    "messages": messages_for_llm,
                    "stream": True,
                }
                if openai_tools:
                    create_kwargs["tools"] = openai_tools

                response = await client.chat.completions.create(**create_kwargs)

                # Accumulate streaming chunks — we need to detect both
                # content tokens AND tool_call deltas in the same stream.
                round_content_chunks: list[str] = []
                # tool_calls_by_index tracks partial tool call arguments
                tool_calls_by_index: dict[int, dict] = {}

                async for chunk in response:
                    choice = chunk.choices[0] if chunk.choices else None
                    if not choice:
                        continue

                    delta = choice.delta

                    # Stream text content to the frontend
                    if delta.content:
                        round_content_chunks.append(delta.content)
                        collected_chunks.append(delta.content)
                        now = time.monotonic()
                        ping = _sse_keepalive(_last_yield, now)
                        if ping:
                            yield ping
                        _last_yield = now
                        yield json.dumps({"type": "token", "content": delta.content})

                    # Accumulate tool calls from streaming deltas
                    if delta.tool_calls:
                        for tc_delta in delta.tool_calls:
                            idx = tc_delta.index
                            if idx not in tool_calls_by_index:
                                tool_calls_by_index[idx] = {
                                    "id": "",
                                    "function": {"name": "", "arguments": ""},
                                }
                            tc = tool_calls_by_index[idx]
                            if tc_delta.id:
                                tc["id"] = tc_delta.id
                            if tc_delta.function:
                                if tc_delta.function.name:
                                    tc["function"]["name"] = tc_delta.function.name
                                if tc_delta.function.arguments:
                                    tc["function"]["arguments"] += tc_delta.function.arguments

                    # Capture usage from streaming chunks (if provider includes it)
                    chunk_usage = getattr(chunk, "usage", None)
                    if chunk_usage and isinstance(getattr(chunk_usage, "prompt_tokens", None), int):
                        accumulated_prompt_tokens += chunk_usage.prompt_tokens or 0
                        accumulated_completion_tokens += chunk_usage.completion_tokens or 0

                    # Detect finish_reason
                    if choice.finish_reason == "tool_calls":
                        break

                # ── Process tool calls ──────────────────────────────
                if tool_calls_by_index:
                    # Build the assistant message with tool_calls for the history
                    assistant_tool_calls = []
                    for idx in sorted(tool_calls_by_index.keys()):
                        tc = tool_calls_by_index[idx]
                        assistant_tool_calls.append(
                            {
                                "id": tc["id"],
                                "type": "function",
                                "function": {
                                    "name": tc["function"]["name"],
                                    "arguments": tc["function"]["arguments"],
                                },
                            }
                        )

                    # Add assistant message with tool_calls to history
                    messages_for_llm.append(
                        {
                            "role": "assistant",
                            "content": "".join(round_content_chunks) or None,
                            "tool_calls": assistant_tool_calls,
                        }
                    )

                    # Execute each tool and add results to history
                    for tc in assistant_tool_calls:
                        tool_name = tc["function"]["name"]
                        tool_args = tc["function"]["arguments"]
                        call_id = tc["id"]

                        yield json.dumps(
                            {
                                "type": "tool_call_start",
                                "tool": tool_name,
                                "arguments": tool_args,
                                "call_id": call_id,
                            }
                        )

                        _tool_t0 = time.time()
                        tool_result = await _execute_tool_call(
                            tool_name,
                            tool_args,
                            user_id=user_id,
                            _user_scopes=_cached_user_scopes,
                            _user_role=_cached_user_role,
                        )
                        _record_tool_cost_fire_and_forget(
                            user_id=user_id,
                            tool_name=tool_name,
                            duration_ms=(time.time() - _tool_t0) * 1000.0,
                            workspace_id=_workspace_id,
                        )

                        yield json.dumps(
                            {
                                "type": "tool_call_result",
                                "tool": tool_name,
                                "result": tool_result,
                                "call_id": call_id,
                            }
                        )

                        # ── Phase 4: canvas_update for agent-driven tile orchestration ──
                        # When a tool produces a result that should open a canvas tile
                        # (e.g. browser_sandbox launch), emit a canvas_update event so
                        # the frontend can auto-open the tile without user action.
                        _canvas_event = _build_canvas_update(tool_name, tool_result)
                        if _canvas_event:
                            yield json.dumps(_canvas_event)

                        messages_for_llm.append(
                            {
                                "role": "tool",
                                "tool_call_id": call_id,
                                "content": tool_result,
                            }
                        )

                    # Task 1.2a: keepalive after tool execution (the longest idle gap)
                    now = time.monotonic()
                    ping = _sse_keepalive(_last_yield, now)
                    if ping:
                        yield ping
                        _last_yield = now

                    # Loop back for the next LLM call with tool results
                    continue

                # No tool calls — we have a final text response
                full_response = "".join(round_content_chunks) or "".join(collected_chunks)
                break  # Exit the tool-calling loop

            else:
                # Loop exhausted — all rounds were tool calls
                full_response = "".join(collected_chunks) or (
                    "I reached the maximum number of tool calls "
                    f"({_MAX_TOOL_ROUNDS}) without producing a final response."
                )

            # Fallback: some providers return 0 streaming chunks even
            # though non-streaming works fine.
            if not full_response.strip():
                logger.info(
                    "stream_message_to_llm: empty stream for %s, retrying without streaming",
                    raw_model,
                )
                try:
                    non_stream_kwargs: dict = {
                        "model": model,
                        "messages": messages_for_llm,
                    }
                    if openai_tools:
                        non_stream_kwargs["tools"] = openai_tools
                    non_stream_response = await client.chat.completions.create(**non_stream_kwargs)
                    if non_stream_response.choices:
                        full_response = non_stream_response.choices[0].message.content or ""
                    # Second fallback: retry without tools if still empty
                    # (some models don't support function calling)
                    if not full_response.strip() and openai_tools:
                        logger.info(
                            "stream_message_to_llm: empty response with tools for %s, retrying without tools",
                            raw_model,
                        )
                        non_stream_kwargs.pop("tools", None)
                        no_tools_response = await client.chat.completions.create(**non_stream_kwargs)
                        if no_tools_response.choices:
                            full_response = no_tools_response.choices[0].message.content or ""
                    if full_response:
                        yield json.dumps({"type": "token", "content": full_response})
                except Exception as retry_err:
                    logger.error(
                        "stream_message_to_llm: non-streaming retry also failed for %s: %s",
                        raw_model,
                        retry_err,
                    )

        # Use actual token counts if available from streaming, else estimate
        prompt_tokens = accumulated_prompt_tokens or len(content.split())
        completion_tokens = accumulated_completion_tokens or len(full_response.split())
        total_tokens = prompt_tokens + completion_tokens

        llm_duration = time.time() - llm_start
        record_llm_request(
            provider=provider_name,
            duration_seconds=llm_duration,
            prompt_tokens=prompt_tokens,
            completion_tokens=completion_tokens,
            success=True,
        )

        # Session was closed before LLM call to prevent idle-in-transaction timeout.
        # Always use fresh session for saving.
        import asyncio

        assistant_msg: ChatMessage | None = None
        # Save with retry (3 attempts, exponential backoff)
        assistant_msg = None
        for _attempt in range(3):
            try:
                assistant_msg = await create_chat_message_fresh_session(thread_id, "assistant", full_response)
                break
            except Exception as save_err:
                logger.warning("stream: assistant message save attempt %d failed: %s", _attempt + 1, save_err)
                if _attempt < 2:
                    await asyncio.sleep(0.5 * (2**_attempt))
        if assistant_msg is None:
            logger.error("stream: assistant message save failed after 3 attempts")
            save_failed_payload = json.dumps({"type": "save_failed", "content": full_response[:500]})
            # Phase 2c.3: emit bare JSON — the outer SSE wrapper frames it
            # with "data: ". Prefixing here double-frames the event.
            yield save_failed_payload

        # ── T33 Stage 1: emit memory_recall_used events for cited claims ──
        # Emitted AFTER the assistant message is persisted so the frontend
        # can attach the recall metadata to the right message. Skipped when
        # the save failed because there's no message_id to attach to.
        if assistant_msg is not None:
            for claim in memory_recall_claims:
                yield json.dumps(build_recall_used_event(claim, message_id=str(assistant_msg.id)))

        # ── T33 Stage 2: emit memory_citation events for the chip UI ──
        # Same ordered list as the recall-used events above. The frontend
        # renders one <MemoryCitationChip> per memory_citation event.
        # The bracket label is built in the service (format_citation_label)
        # so the chip displays the locked format verbatim.
        if assistant_msg is not None:
            for claim in memory_recall_claims:
                yield json.dumps(build_citation_event(claim, message_id=str(assistant_msg.id)))

        # ── P0-1: durable memory extraction via Celery ────────────────
        try:
            from app.tasks.memory_extraction_tasks import extract_memory_claims_task

            extract_memory_claims_task.delay(
                thread_id=thread_id,
                user_id=user_id,
                user_message=content,
                assistant_response=full_response,
            )
        except Exception:
            logger.debug("memory_extraction Celery dispatch failed (non-fatal)", exc_info=True)

        try:
            from app.services.usage_service import get_usage_service

            get_usage_service().record_usage(
                user_id=str(user_id),
                model_id=model,
                provider="byok" if user_api_key else "system",
                prompt_tokens=prompt_tokens,
                completion_tokens=completion_tokens,
                cost=total_tokens * 0.000002,
            )
        except Exception as ue:
            logger.warning("Usage recording failed (non-fatal): %s", ue)

        yield json.dumps(
            {
                "type": "complete",
                "full_response": full_response,
                "message_id": assistant_msg.id if assistant_msg is not None else None,
                "model": model,
            }
        )

    except CircuitOpenError as e:
        llm_duration = time.time() - llm_start
        record_llm_request(provider=provider_name, duration_seconds=llm_duration, success=False)
        logger.warning("Circuit breaker open for %s: %s", provider_name, e)
        yield json.dumps(
            {
                "type": "error",
                "error": f"Service temporarily unavailable ({provider_name}). Please try again later.",
            }
        )
    except Exception as e:
        llm_duration = time.time() - llm_start
        record_llm_request(provider=provider_name, duration_seconds=llm_duration, success=False)
        logger.error("stream_message_to_llm failed: %s", e)
        yield json.dumps({"type": "error", "error": str(e)})


def _sse_keepalive(last_yield_time: float, now: float) -> str | None:
    """Return a SSE keepalive comment if enough time has elapsed since the last yield."""
    if now - last_yield_time > _SSE_KEEPALIVE_INTERVAL:
        return ": ping\n\n"
    return None


_client = AsyncOpenAI(
    api_key=_LLM_API_KEY,
    base_url=_LLM_API_BASE,
)


async def create_chat_thread(
    db: AsyncSession,
    user_id: int,
    username: str,
    title: str,
    model_preference: str | None = None,
    workspace_id: str | None = None,
) -> ChatThread:
    metadata = {"model_preference": model_preference} if model_preference else None
    thread = ChatThread(
        title=title,
        user_id=user_id,
        username=username,
        metadata_=metadata,
        workspace_id=workspace_id,
    )
    db.add(thread)
    await db.flush()
    await db.refresh(thread)
    return thread


async def get_chat_thread(db: AsyncSession, thread_id: int) -> ChatThread | None:
    result = await db.execute(select(ChatThread).where(ChatThread.id == thread_id))
    return result.scalar_one_or_none()


async def require_chat_thread_access(
    db: AsyncSession,
    thread_id: int,
    user_id: int,
) -> ChatThread:
    """Fetch a chat thread and verify the user has access.

    Access rules:
    1. If the thread has a workspace_id → verify the user is an active member
       of that workspace.
    2. If the thread has no workspace_id → fall back to user_id ownership.
    3. If the thread doesn't exist → 404.
    """
    from fastapi import HTTPException

    thread = await get_chat_thread(db, thread_id)
    if thread is None:
        raise HTTPException(status_code=404, detail="Not found")

    if thread.workspace_id:
        from app.models.workspace_models import WorkspaceMember

        result = await db.execute(
            select(WorkspaceMember).where(
                WorkspaceMember.workspace_id == thread.workspace_id,
                WorkspaceMember.user_id == user_id,
                WorkspaceMember.is_active == True,
            )
        )
        if result.scalar_one_or_none() is None:
            # Check cross-workspace grants across all user's workspaces
            from app.services.cross_workspace_service import (
                check_entity_access,
                find_user_workspaces,
            )

            user_workspaces = await find_user_workspaces(db, user_id)
            has_cross_access = False
            for ws_id in user_workspaces:
                if ws_id == thread.workspace_id:
                    continue
                grant = await check_entity_access(
                    db,
                    user_id=user_id,
                    target_workspace_id=ws_id,
                    entity_type="chat_thread",
                    entity_id=str(thread_id),
                    required_permission="read",
                )
                if grant:
                    has_cross_access = True
                    break
            if not has_cross_access:
                logger.warning(
                    "entity_access_denied"
                    " user_id=%s entity_type=chat_thread entity_id=%s"
                    " workspace_id=%s reason=no_membership",
                    user_id,
                    thread_id,
                    thread.workspace_id,
                )
                try:
                    from app.api.middleware.audit import log_event
                    from app.services.background_task_manager import background_task_manager

                    background_task_manager.spawn(
                        log_event(
                            user_id=user_id,
                            action="entity.access_denied",
                            details={
                                "entity_type": "chat_thread",
                                "entity_id": str(thread_id),
                                "workspace_id": str(thread.workspace_id),
                                "reason": "no_membership",
                            },
                        ),
                        label="access_denied_audit_no_membership",
                    )
                except Exception:
                    pass
                raise HTTPException(status_code=404, detail="Not found")
    else:
        if thread.user_id != user_id:
            logger.warning(
                "entity_access_denied"
                " user_id=%s entity_type=chat_thread entity_id=%s"
                " owner_user_id=%s reason=owner_mismatch",
                user_id,
                thread_id,
                thread.user_id,
            )
            try:
                from app.api.middleware.audit import log_event
                from app.services.background_task_manager import background_task_manager

                background_task_manager.spawn(
                    log_event(
                        user_id=user_id,
                        action="entity.access_denied",
                        details={
                            "entity_type": "chat_thread",
                            "entity_id": str(thread_id),
                            "reason": "owner_mismatch",
                        },
                    ),
                    label="access_denied_audit_owner_mismatch",
                )
            except Exception:
                pass
            raise HTTPException(status_code=404, detail="Not found")

    return thread


async def list_chat_threads(
    db: AsyncSession,
    user_id: int,
    *,
    offset: int = 0,
    limit: int = 20,
    workspace_id: str | None = None,
) -> tuple[list[ChatThread], int]:
    base_filter = ChatThread.workspace_id == workspace_id if workspace_id is not None else ChatThread.user_id == user_id
    count_q = select(func.count()).select_from(ChatThread).where(base_filter)
    total = (await db.execute(count_q)).scalar() or 0
    items_q = select(ChatThread).where(base_filter).order_by(ChatThread.updated_at.desc()).offset(offset).limit(limit)
    items = list((await db.execute(items_q)).scalars().all())
    return items, total


async def update_chat_thread(
    db: AsyncSession,
    thread_id: int,
    *,
    title: str | None = None,
    is_archived: bool | None = None,
) -> ChatThread | None:
    thread = await get_chat_thread(db, thread_id)
    if thread is None:
        return None
    if title is not None:
        thread.title = title
    if is_archived is not None:
        thread.is_archived = is_archived
    await db.flush()
    await db.refresh(thread)
    return thread


async def delete_chat_thread(db: AsyncSession, thread_id: int) -> bool:
    thread = await get_chat_thread(db, thread_id)
    if thread is None:
        return False
    await db.delete(thread)
    await db.flush()
    return True


async def create_chat_message(
    db: AsyncSession,
    thread_id: int,
    role: str,
    content: str,
    *,
    user_id: int | None = None,
) -> ChatMessage:
    msg = ChatMessage(thread_id=thread_id, role=role, content=content)
    if user_id is not None:
        msg.user_id = user_id
    db.add(msg)
    await db.flush()
    await db.refresh(msg)
    return msg


async def create_chat_message_fresh_session(
    thread_id: int,
    role: str,
    content: str,
    *,
    user_id: int | None = None,
) -> ChatMessage:
    """Create a chat message using a fresh DB session.

    Used as a fallback when the caller's session has a dead connection
    (e.g. after long LLM streaming where idle-in-transaction kills the
    underlying asyncpg connection).  The session is committed by the
    context manager on successful exit.
    """
    from app.database import fresh_session

    async with fresh_session() as fresh_db:
        msg = ChatMessage(thread_id=thread_id, role=role, content=content)
        if user_id is not None:
            msg.user_id = user_id
        fresh_db.add(msg)
        await fresh_db.flush()
        await fresh_db.refresh(msg)
        return msg


async def update_chat_message(
    db: AsyncSession,
    message_id: int,
    content: str,
) -> ChatMessage | None:
    """Update a chat message's content and set edited_at timestamp."""
    from datetime import datetime

    result = await db.execute(select(ChatMessage).where(ChatMessage.id == message_id))
    message = result.scalar_one_or_none()
    if message is None:
        return None
    message.content = content
    message.edited_at = datetime.now(UTC)
    await db.flush()
    await db.refresh(message)
    return message


async def delete_chat_message(db: AsyncSession, message_id: int) -> bool:
    """Delete a chat message by ID."""
    result = await db.execute(select(ChatMessage).where(ChatMessage.id == message_id))
    message = result.scalar_one_or_none()
    if message is None:
        return False
    await db.delete(message)
    await db.flush()
    return True


async def get_chat_messages(
    db: AsyncSession,
    thread_id: int,
    *,
    offset: int = 0,
    limit: int = 50,
) -> tuple[list[ChatMessage], int]:
    count_q = select(func.count()).select_from(ChatMessage).where(ChatMessage.thread_id == thread_id)
    total = (await db.execute(count_q)).scalar() or 0
    items_q = (
        select(ChatMessage)
        .where(ChatMessage.thread_id == thread_id)
        .order_by(ChatMessage.created_at.asc())
        .offset(offset)
        .limit(limit)
    )
    items = list((await db.execute(items_q)).scalars().all())
    return items, total


async def get_chat_files(
    db: AsyncSession,
    thread_id: int,
) -> list[ChatFile]:
    result = await db.execute(select(ChatFile).where(ChatFile.chat_id == thread_id))
    return list(result.scalars().all())


async def create_chat_file(
    db: AsyncSession,
    thread_id: int,
    filename: str,
    mime_type: str,
    path: str,
    size_bytes: int,
) -> ChatFile:
    file = ChatFile(
        chat_id=thread_id,
        filename=filename,
        mime_type=mime_type,
        path=path,
        size_bytes=size_bytes,
    )
    db.add(file)
    await db.flush()
    await db.refresh(file)
    return file


def _get_model_preference(thread: ChatThread) -> str | None:
    if thread.metadata_ and isinstance(thread.metadata_, dict):
        return thread.metadata_.get("model_preference")
    return None


async def _process_attachments(
    db: AsyncSession, messages: list[dict], attachments: list[dict], model: str
) -> list[dict]:
    is_vision_model = not model.startswith("llamacpp/")

    for att in attachments:
        file_id = att.get("file_id", "")
        att_type = att.get("type", "file")
        filename = att.get("filename", "unknown")

        result = await db.execute(select(UserFile).where(UserFile.id == file_id))
        db_file = result.scalar_one_or_none()
        if not db_file or not db_file.storage_path or not os.path.exists(db_file.storage_path):
            continue

        if att_type == "image" and is_vision_model:
            raw_bytes = Path(db_file.storage_path).read_bytes()
            b64 = base64.b64encode(raw_bytes).decode("utf-8")
            content_type = db_file.content_type or "image/png"

            if messages and messages[-1].get("role") == "user":
                existing = messages[-1].get("content", "")
                if isinstance(existing, str):
                    messages[-1]["content"] = [
                        {"type": "text", "text": existing},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{content_type};base64,{b64}"},
                        },
                    ]
                else:
                    messages[-1]["content"].append(
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{content_type};base64,{b64}"},
                        }
                    )

        elif att_type == "image" and not is_vision_model:
            if messages and messages[-1].get("role") == "user":
                existing = messages[-1].get("content", "")
                if isinstance(existing, str):
                    messages[-1]["content"] = f"{existing}\n\n[Image: {filename}]"

        elif att_type == "file":
            try:
                content_type = (db_file.content_type or "").lower()
                filename_lower = (db_file.filename or "").lower()

                if "pdf" in content_type:
                    import pdfplumber

                    with pdfplumber.open(db_file.storage_path) as pdf:
                        pages_text = []
                        for page in pdf.pages:
                            t = page.extract_text() or ""
                            pages_text.append(t)
                        file_text = "\n\n".join(pages_text)

                elif "wordprocessingml" in content_type or ".docx" in filename_lower or "msword" in content_type:
                    from docx import Document

                    doc = Document(db_file.storage_path)
                    file_text = "\n".join(p.text for p in doc.paragraphs)

                elif "spreadsheetml" in content_type or ".xlsx" in filename_lower or "excel" in content_type:
                    import openpyxl

                    wb = openpyxl.load_workbook(db_file.storage_path, read_only=True, data_only=True)
                    rows = []
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        sheet_rows = []
                        for row in ws.iter_rows(values_only=True):
                            line = "\t".join(str(c) if c is not None else "" for c in row)
                            sheet_rows.append(line)
                        rows.append(f"=== Sheet: {sheet} ===\n" + "\n".join(sheet_rows))
                    file_text = "\n\n".join(rows)
                    wb.close()

                elif "presentationml" in content_type or ".pptx" in filename_lower or "powerpoint" in content_type:
                    from pptx import Presentation

                    prs = Presentation(db_file.storage_path)
                    slides_text = []
                    for i, slide in enumerate(prs.slides, 1):
                        slide_texts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                slide_texts.append(shape.text)
                            elif shape.has_table:
                                table = shape.table
                                slide_texts.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
                        slides_text.append(f"=== Slide {i} ===\n" + "\n".join(slide_texts))
                    file_text = "\n\n".join(slides_text)

                else:
                    file_text = Path(db_file.storage_path).read_text(encoding="utf-8", errors="replace")
            except Exception:
                continue
            limit = 10000
            truncated = file_text[:limit]
            if len(file_text) > limit:
                truncated += "\n... (truncated)"
            context_msg = f"[Attached file: {filename}]\n{truncated}\n[End of attached file]"
            messages.insert(-1, {"role": "user", "content": context_msg})

    return messages


async def _prepare_step_inject(
    messages: list[dict],
    *,
    memory_claims: list | None = None,
    web_search: bool = False,
    content: str | None = None,
    steps: list | None = None,
    should_inject: bool | None = None,
) -> tuple[list[dict], list[dict]]:
    """Ordered context-injection closure for chat (SPIKE — ADR-002).

    Mirrors trigger.dev ``chat.agent``'s single ``prepareStep`` primitive:
    all context sources are merged into ``messages`` in a fixed order
    (memory -> web search) and the function returns both the mutated message
    list and a list of receipt events.

    Today this is called once per turn, pre-LLM, with ``steps=None``
    (single-shot chat). The future re-entrant turn loop will call it at each
    AI-SDK step boundary with a non-empty ``steps`` list and a ``should_inject``
    gate, matching trigger.dev's step-boundary injection model. The ordering is
    enforced here by code structure, not convention.

    Behavior when the flag is off is identical to the legacy inline path
    (see ``_stream_message_to_llm_body``); this helper is only reached when
    ``settings.CHAT_PREPARE_STEP_HOOK_ENABLED`` is True.

    Returns:
        (messages_with_context, injected_events) where each event is an SSE
        payload dict the caller yields so the frontend can reconcile
        injected-vs-queued context.
    """
    injected_events: list[dict] = []
    result = list(messages)

    # 1) Memory context (ordered first, like prepareStep's compaction->steering)
    if memory_claims:
        result = _inject_memory_context(result, memory_claims)
        injected_events.append(
            {
                "type": "injected",
                "source": "memory",
                "count": len(memory_claims),
            }
        )

    # 2) Web search context
    if web_search and content is not None:
        result = await _inject_web_search(result, content)
        injected_events.append(
            {
                "type": "injected",
                "source": "web_search",
                "query": content,
            }
        )

    # (Future step-boundary hook point) if steps is not None, apply a
    # should_inject gate and additional per-step sources here.
    return result, injected_events


async def _inject_web_search(messages: list[dict], query: str) -> list[dict]:
    try:
        from app.services.web_search.models import SearchConfig
        from app.services.web_search.service_enhanced import EnhancedWebSearchService

        config = SearchConfig(
            duckduckgo_enabled=True,
            searxng_enabled=True,
            searxng_url="http://searxng:8080",
        )
        service = EnhancedWebSearchService(config)
        # Extract key search terms for better results
        search_query = query
        for prefix in [
            "what is",
            "what's",
            "what are",
            "tell me",
            "can you",
            "do you know",
            "how much is",
            "how much does",
            "current",
            "latest",
            "find",
            "search",
        ]:
            if search_query.lower().startswith(prefix):
                search_query = search_query[len(prefix) :].strip(" ?!,.;:")
                break
        if not search_query.strip():
            search_query = query
        search_result = await service.search(
            query=search_query,
            max_results=5,
            use_cache=False,
            use_reranking=False,
            use_query_understanding=False,
        )

        results = search_result.get("results", [])
        if not results:
            return messages

        lines = [
            "IMPORTANT: The following web search results contain CURRENT, FACTUAL information.",
            f'Use these results to answer the user\'s question about: "{query}"',
            "Do NOT say you lack access to live data — you have the data right here.",
            "If the results contain a price, date, or figure, cite it directly.",
            "",
            "Search results:",
        ]
        for i, r in enumerate(results[:5], 1):
            if isinstance(r, dict):
                title = r.get("title", "Untitled")
                url = r.get("url", "")
                snippet = r.get("snippet", r.get("content", ""))[:200]
            else:
                title = getattr(r, "title", "Untitled") or "Untitled"
                url = getattr(r, "url", "") or ""
                snippet = (getattr(r, "snippet", "") or getattr(r, "content", "") or "")[:200]
            lines.append(f"{i}. {title} — {url}")
            if snippet:
                lines.append(f"   {snippet}")
            lines.append("")

        search_context = "\n".join(lines)
        # Append search results to the last user message so the model can't ignore them
        messages[-1]["content"] = (
            f"{search_context}\n\nBased on the search results above, answer the following:\n{messages[-1]['content']}"
        )

    except Exception as e:
        logger.warning("Web search failed (non-fatal): %s", e)

    return messages


# ── Phase 6: prompt version caching ───────────────────────────────
# Same Redis caching pattern as workspace tool allowlist (5min TTL).
# Avoids a DB query on every chat message.
_PROMPT_CACHE_TTL = 300  # 5 minutes


def _prompt_cache_key(workspace_id: str, name: str) -> str:
    return f"prompt_version:{workspace_id}:{name}"


async def _get_prompt_redis():
    """Return an async Redis client, or None if unavailable."""
    try:
        from redis.asyncio import from_url

        url = os.getenv("REDIS_URL", "redis://localhost:6379/0")
        return from_url(url, decode_responses=True)
    except Exception:
        return None


async def invalidate_prompt_version_cache(workspace_id: str, name: str | None = None) -> None:
    """Delete cached prompt version(s) for a workspace.

    Called after create/activate/delete in the prompt CRUD API.
    If *name* is given, invalidates only that name group.
    If *name* is None, invalidates all prompt caches for the workspace.
    """
    rds = await _get_prompt_redis()
    if rds is None:
        return
    try:
        if name:
            await rds.delete(_prompt_cache_key(workspace_id, name))
        else:
            # Scan for all prompt_version:{workspace_id}:* keys
            pattern = f"prompt_version:{workspace_id}:*"
            cursor = 0
            while True:
                cursor, keys = await rds.scan(cursor, match=pattern, count=100)
                if keys:
                    await rds.delete(*keys)
                if cursor == 0:
                    break
    except Exception:
        pass  # cache miss is fine
    finally:
        await rds.aclose()


async def _get_active_prompt_content(
    db: AsyncSession,
    workspace_id: str,
    name: str = "Default Assistant",
) -> str | None:
    """Look up the active prompt version for a workspace and name group.

    Returns the prompt content string, or None if no active version exists.
    Phase 6: prompt versioning with Redis caching (5min TTL).

    Results are cached in Redis to avoid a DB query on every chat message.
    Cache is invalidated by ``invalidate_prompt_version_cache()`` when
    prompt versions are created, activated, or deleted via the API.
    """
    from app.models.prompt_version_models import PromptVersion

    cache_key = _prompt_cache_key(workspace_id, name)

    # ── Single Redis connection for read + write ───────────────────
    rds = await _get_prompt_redis()
    try:
        # ── Try cache first ──────────────────────────────────────
        if rds is not None:
            try:
                cached = await rds.get(cache_key)
                if cached is not None:
                    if cached == "__NONE__":
                        return None  # sentinel: no active version
                    return cached
            except Exception:
                pass  # cache miss → fall through to DB

        # ── DB query ───────────────────────────────────────────
        result = await db.execute(
            select(PromptVersion).where(
                PromptVersion.workspace_id == workspace_id,
                PromptVersion.name == name,
                PromptVersion.is_active == True,
            )
        )
        pv = result.scalar_one_or_none()
        content = pv.content if pv is not None else None

        # ── Populate cache (same connection) ───────────────────
        if rds is not None:
            try:
                cache_val = content if content is not None else "__NONE__"
                await rds.setex(cache_key, _PROMPT_CACHE_TTL, cache_val)
            except Exception:
                pass  # caching is best-effort

        return content
    finally:
        if rds is not None:
            await rds.aclose()


async def _build_chat_messages(
    db: AsyncSession,
    thread_id: int,
    max_history: int = 20,
) -> list[dict]:
    # Build the messages array for the LLM including conversation history.
    # Fetches the last max_history messages from the thread and includes them
    # in the prompt. Uses the thread custom system prompt if available.
    #
    # Phase 6: prompt version lookup chain:
    #   1. If thread has a workspace_id, try prompt_versions table first
    #   2. Fall back to thread.metadata_.get("system_prompt")
    #   3. Final fallback: "You are a helpful assistant."
    thread = await get_chat_thread(db, thread_id)

    system_prompt: str | None = None

    # Phase 6: try prompt_versions first (workspace-scoped)
    if thread and thread.workspace_id:
        # Check if thread has a specific prompt_version_id stored
        prompt_name = "Default Assistant"
        if thread.metadata_ and isinstance(thread.metadata_, dict):
            prompt_name = thread.metadata_.get("prompt_name", "Default Assistant")
        system_prompt = await _get_active_prompt_content(db, thread.workspace_id, name=prompt_name)

    # Fallback: inline system prompt from thread metadata
    if not system_prompt and thread and thread.metadata_ and isinstance(thread.metadata_, dict):
        custom = thread.metadata_.get("system_prompt", "")
        if custom and custom.strip():
            system_prompt = custom

    if not system_prompt:
        system_prompt = "You are a helpful assistant."

    # Append sandboxd preview guidance when sandboxd is enabled
    if settings.SANDBOXD_ENABLED:
        system_prompt += _SANDBOXD_SYSTEM_GUIDANCE

    messages = [{"role": "system", "content": system_prompt}]

    history_stmt = (
        select(ChatMessage).where(ChatMessage.thread_id == thread_id).order_by(ChatMessage.id.desc()).limit(max_history)
    )
    history_result = await db.execute(history_stmt)
    recent_messages = list(reversed(history_result.scalars().all()))

    messages.extend(
        {"role": msg.role, "content": msg.content} for msg in recent_messages if msg.role in ("user", "assistant")
    )

    # Phase 2.1: token-budget pruning (keeps first+last, replaces middle)
    if settings.CHAT_CONTEXT_PRUNING_ENABLED:
        messages = _prune_messages_to_budget(messages, settings.CHAT_CONTEXT_TOKEN_BUDGET)

    return messages


async def _maybe_extract_memory_claims(
    *,
    db: AsyncSession | None = None,
    thread_id: int,
    user_id: int,
    user_message: str,
    assistant_response: str,
) -> None:
    """Fire-and-forget: extract personal-memory claims from a chat exchange.

    Called after the assistant response is persisted.  Opens a *fresh*
    DB session so it can run independently of the caller's session
    lifecycle (the caller's session is typically closed before the
    LLM call to avoid idle-in-transaction timeout).

    Guardrails:
    * Gated by ``FLOWMANNER_CROSS_MISSION_MEMORY`` feature flag.
    * Skipped when the conversation is paused.
    * Defensive filter drops sensitive/restricted/private claims.
    * Capped at 5 claims per exchange.
    * Errors are swallowed — never breaks the chat.
    """
    if not settings.FLOWMANNER_CROSS_MISSION_MEMORY:
        return

    import asyncio

    try:
        from app.database import fresh_session

        async with fresh_session() as fresh_db:
            # Need workspace_id from the thread.
            thread = await get_chat_thread(fresh_db, thread_id)
            if thread is None or not thread.workspace_id:
                return

            workspace_id = thread.workspace_id

            # ── Check pause toggle ──────────────────────────────────
            from app.services.memory_extraction_pause_service import (
                MemoryExtractionPauseService,
            )

            pause_svc = MemoryExtractionPauseService(fresh_db)
            if await pause_svc.is_paused(
                user_id=user_id,
                workspace_id=workspace_id,
                conversation_id=str(thread_id),
            ):
                logger.info(
                    "memory_extraction: skipped for thread %s (paused)",
                    thread_id,
                )
                return

            # ── Run extractor ───────────────────────────────────────
            # Attempt LLM-based extraction first (via ModelRouter +
            # PersonalMemoryExtractor) with a 5-second timeout.  If the
            # LLM is unavailable, slow, or returns garbage, fall back
            # to the deterministic regex extractor.
            combined_text = f"User: {user_message}\n\nAssistant: {assistant_response}"

            regex_fallback = RegexPersonalMemoryExtractor()
            claims: list = []
            _extraction_source = "empty"
            try:
                from app.services.model_router import get_model_router

                llm_extractor = PersonalMemoryExtractor(
                    get_model_router=get_model_router,
                )
                claims, source = await asyncio.wait_for(
                    llm_extractor.extract_with_fallback(
                        user_id=user_id,
                        workspace_id=workspace_id,
                        text=combined_text,
                        max_claims=5,
                    ),
                    timeout=5.0,
                )
                if claims:
                    _extraction_source = "llm"
                    logger.info(
                        "memory_extraction: LLM extractor returned %d claims (source=%s) for thread %s",
                        len(claims),
                        source,
                        thread_id,
                    )
                else:
                    # LLM succeeded but produced nothing — try regex
                    claims = regex_fallback.extract(combined_text)
                    _extraction_source = "regex_fallback_empty"
                    if claims:
                        logger.info(
                            "memory_extraction: LLM returned 0 claims; regex fallback found %d for thread %s",
                            len(claims),
                            thread_id,
                        )
            except TimeoutError:
                logger.info(
                    "memory_extraction: LLM extraction timed out for thread %s; falling back to regex",
                    thread_id,
                )
                claims = regex_fallback.extract(combined_text)
                _extraction_source = "regex_fallback_timeout"
            except Exception as llm_err:
                logger.info(
                    "memory_extraction: LLM extraction failed for thread %s (%s); falling back to regex",
                    thread_id,
                    llm_err,
                )
                claims = regex_fallback.extract(combined_text)
                _extraction_source = "regex_fallback_error"

            if not claims:
                logger.debug(
                    "memory_extraction: no claims extracted for thread %s",
                    thread_id,
                )
                return

            # ── Defensive filter ────────────────────────────────────
            # Drop sensitive/restricted claims and private-scope claims.
            # CandidateClaim has no ``sensitivity`` field (it is a
            # lightweight DTO), so we fall back to ``claim_type`` as a
            # belt-and-suspenders check.
            _EXCLUDED_SENSITIVITIES = frozenset({"sensitive", "restricted"})
            _EXCLUDED_SCOPES = frozenset({"private"})

            safe_claims = [
                c
                for c in claims
                if getattr(c, "sensitivity", "normal") not in _EXCLUDED_SENSITIVITIES
                and c.claim_type not in _EXCLUDED_SENSITIVITIES
                and c.scope not in _EXCLUDED_SCOPES
            ]

            # GOV-1.5 (C5): persist *why* claims were dropped. The
            # defensive filter previously dropped claims silently, so the
            # 0.85 confidence gate could never be calibrated from real
            # data. Log every defensively-dropped claim with its score
            # and reason so a later pass can quantify the drop rate.
            dropped_defensive = len(claims) - len(safe_claims)
            if dropped_defensive:
                _dropped_below: list[float] = []
                for c in claims:
                    if c not in safe_claims:
                        _dropped_below.append(float(getattr(c, "confidence", 0.0)))
                        logger.info(
                            "memory_extraction: dropped (defensive filter) "
                            "thread=%s claim_type=%s scope=%s confidence=%.2f "
                            "subject=%s predicate=%s",
                            thread_id,
                            getattr(c, "claim_type", "?"),
                            getattr(c, "scope", "?"),
                            float(getattr(c, "confidence", 0.0)),
                            getattr(c, "subject", "?"),
                            getattr(c, "predicate", "?"),
                        )
                logger.info(
                    "memory_extraction: defensive filter dropped %d/%d claims for thread %s (scores=%s)",
                    dropped_defensive,
                    len(claims),
                    thread_id,
                    _dropped_below,
                )

            if not safe_claims:
                logger.info(
                    "memory_extraction: all %d claims filtered for thread %s",
                    len(claims),
                    thread_id,
                )
                return

            # ── Persist claims ──────────────────────────────────────
            # Solo workspaces → direct write (no approval needed).
            # Team workspaces (multi-member, >30 days) → stage for
            # user approval via BackgroundReviewService.
            from app.models.workspace_models import Workspace
            from app.services.memory.background_review_service import (
                BackgroundReviewService,
                compute_write_approval,
            )
            from app.services.personal_memory_service import PersonalMemoryService

            # Fetch workspace to determine approval path.
            ws_row = (
                await fresh_db.execute(select(Workspace).where(Workspace.id == workspace_id))
            ).scalar_one_or_none()
            needs_approval = compute_write_approval(ws_row)

            from app.services.memory.extraction_thresholds import (
                MEMORY_EXTRACTION_MIN_CONFIDENCE,
                is_trusted_direct_write,
                passes_confidence_gate,
            )
            from app.services.memory.provenance_approval import (
                requires_provenance_approval,
            )

            pm_service = PersonalMemoryService(fresh_db)
            review_service = BackgroundReviewService()
            persisted = 0
            staged = 0
            for claim in safe_claims:
                # GOV-1.2 provenance gate (deterministic, no confidence
                # bypass). The chat extractor infers every claim here with
                # source_type="conversation" (never user_explicit), so the
                # gate routes ALL extracted claims to human approval —
                # unless the workspace/approval decision already forces
                # staging. Only a literally user-authored source_type
                # ("user_explicit") may bypass to a direct write.
                # This is the reliable control; 1.3a scan / 1.3b scrub are
                # NOT allowed to de-escalate it.
                claim_source_type = getattr(claim, "source_type", None) or "conversation"
                provenance_requires_approval = requires_provenance_approval(claim_source_type)
                try:
                    if needs_approval or provenance_requires_approval:
                        # Stage for user approval
                        content_str = f"{claim.subject} {claim.predicate}: {json.dumps(claim.object, default=str)}"
                        pw_id = await review_service.stage_pending_write(
                            fresh_db,
                            workspace_id=workspace_id,
                            user_id=user_id,
                            mission_id=None,
                            action="add",
                            content=content_str,
                            metadata={"source_type": claim_source_type},
                        )
                        if pw_id:
                            staged += 1
                    else:
                        # Direct write — only reachable for a user_explicit
                        # source_type (see requires_provenance_approval).
                        # GOV-1.5 (calibration): even trusted direct writes
                        # are held for approval when their extractor
                        # confidence is below the calibrated floor. This
                        # gate applies ONLY to the trusted path — untrusted
                        # source_types never reach this branch, so the
                        # confidence gate can never de-escalate a
                        # provenance-mandated approval (GOV-1.2 invariant).
                        if is_trusted_direct_write(claim_source_type) and not passes_confidence_gate(claim.confidence):
                            logger.info(
                                "memory_extraction: trusted claim held for "
                                "approval below confidence gate (%.2f < %.2f) "
                                "thread=%s subject=%s predicate=%s",
                                claim.confidence,
                                MEMORY_EXTRACTION_MIN_CONFIDENCE,
                                thread_id,
                                claim.subject,
                                claim.predicate,
                            )
                            content_str = f"{claim.subject} {claim.predicate}: {json.dumps(claim.object, default=str)}"
                            pw_id = await review_service.stage_pending_write(
                                fresh_db,
                                workspace_id=workspace_id,
                                user_id=user_id,
                                mission_id=None,
                                action="add",
                                content=content_str,
                                metadata={
                                    "source_type": claim_source_type,
                                    "held_reason": "confidence_below_gate",
                                    "confidence": claim.confidence,
                                },
                            )
                            if pw_id:
                                staged += 1
                            continue
                        await pm_service.create(
                            user_id=user_id,
                            workspace_id=workspace_id,
                            subject=claim.subject,
                            predicate=claim.predicate,
                            object=claim.object,
                            claim_type=claim.claim_type,
                            scope=claim.scope,
                            source_type=claim_source_type,
                            confidence=claim.confidence,
                        )
                        persisted += 1
                except Exception as create_err:
                    logger.warning(
                        "memory_extraction: claim persist failed: %s",
                        create_err,
                    )

            # fresh_session() commits on successful exit

            # ── Record extraction metrics ───────────────────────────
            from app.core.metrics import record_memory_extraction

            record_memory_extraction(
                source=_extraction_source,
                claims_extracted=len(claims),
                claims_persisted=persisted,
                claims_staged=staged,
                claims_dropped=dropped_defensive,
            )

            logger.info(
                "memory_extraction: persisted=%d staged=%d/%d raw=%d "
                "claims for thread %s (needs_approval=%s source=%s)",
                persisted,
                staged,
                len(safe_claims),
                len(claims),
                thread_id,
                needs_approval,
                _extraction_source,
            )

            # ── GOV-1.6 (C5): persist dropped candidates durably ───────
            # GOV-1.5 made drops observable (logs + metrics), but the
            # dropped candidates were not durable or Inspector-visible —
            # the GOV-1.6 (C5) gap. Write one ``drop`` MemoryCorrectionEvent
            # per defensively-dropped candidate so the drop rate the 0.85
            # gate calibrates against becomes a queryable signal in the same
            # privacy trail. claim_id is NULL (the candidate never became a
            # PersonalMemoryClaim); its shape is carried in details.
            # No-fail: an audit-sink outage must never break memory capture.
            if dropped_defensive:
                try:
                    from app.services.memory_correction_service import (
                        ALL_EVENT_TYPES,
                        MemoryCorrectionService,
                    )

                    if "drop" in ALL_EVENT_TYPES:
                        drop_svc = MemoryCorrectionService(fresh_db)
                        for c in claims:
                            if c in safe_claims:
                                continue
                            await drop_svc.record_event(
                                user_id=user_id,
                                workspace_id=workspace_id,
                                event_type="drop",
                                claim_id=None,
                                actor="system",
                                source="memory_extraction",
                                details={
                                    "reason": "defensive_filter",
                                    "claim_type": getattr(c, "claim_type", None),
                                    "scope": getattr(c, "scope", None),
                                    "confidence": float(getattr(c, "confidence", 0.0)),
                                    "subject": getattr(c, "subject", None),
                                    "predicate": getattr(c, "predicate", None),
                                },
                            )
                except Exception as drop_err:  # pragma: no cover - no-fail sink
                    logger.warning(
                        "memory_extraction: durable drop audit failed for thread %s: %s",
                        thread_id,
                        drop_err,
                    )
    except Exception as exc:
        logger.warning(
            "memory_extraction: hook failed for thread %s (non-fatal): %s",
            thread_id,
            exc,
        )


def _record_tool_cost_fire_and_forget(
    user_id: int,
    tool_name: str,
    duration_ms: float,
    workspace_id: str | None = None,
) -> None:
    """Fire-and-forget: record tool call cost using a fresh DB session.

    Opens its own ``AsyncSessionLocal`` so the recording is independent of the
    caller's (already-closed) request session.  Errors are swallowed — cost
    tracking must never break the chat.
    """
    import asyncio

    async def _run() -> None:
        try:
            from app.database import fresh_session
            from app.services.cost_tracker import get_cost_tracker

            async with fresh_session() as fresh_db:
                await get_cost_tracker().record_tool_call_cost(
                    db=fresh_db,
                    user_id=user_id,
                    tool_name=tool_name,
                    duration_ms=duration_ms,
                    workspace_id=workspace_id,
                )
                # fresh_session() commits on successful exit
        except Exception as e:
            logger.debug("tool_cost_tracking_failed tool=%s error=%s", tool_name, e)

    from app.services.background_task_manager import background_task_manager

    background_task_manager.spawn(_run(), label="tool_cost_recording")


async def send_message_to_llm(
    db: AsyncSession,
    thread_id: int,
    content: str,
    user_id: int,
    model_preference: str | None = None,
    user_api_key: str | None = None,
    user_base_url: str | None = None,
    model_id: str | None = None,
    attachments: list | None = None,
    web_search: bool | None = None,
) -> dict:
    """Send a message to the LLM and get a non-streaming response.

    If user_api_key is provided, a per-request AsyncOpenAI client is created
    using that key (BYOK path). The key is used only for this call and discarded.
    If model_id is provided it overrides model_preference and the env default.
    If user_base_url is provided, it overrides the resolved base_url for custom providers.
    For llamacpp/* models, any BYOK key is ignored (llama.cpp doesn't use API keys).
    """
    # --- Pre-LLM setup: resolve provider, BYOK lookup, build messages ---
    # All DB work happens BEFORE saving/committing the user message so the
    # session can be closed immediately after commit, preventing PostgreSQL's
    # idle_in_transaction_session_timeout from killing the connection during
    # the potentially multi-minute LLM call.
    raw_model = model_id or model_preference or _LLM_MODEL
    base_url, api_key, model = _resolve_provider(raw_model)

    mismatch_error = _validate_byok_key_matches_model(user_api_key, raw_model)
    if mismatch_error:
        return {
            "success": False,
            "content": mismatch_error,
            "tokens": 0,
            "model": model,
        }

    # --- Stored-key fallback: if no per-request key, look up user's stored key ---
    effective_user_key = user_api_key
    effective_base_url = user_base_url
    if not effective_user_key and db is not None:
        model_provider = _get_provider_for_model(raw_model)
        stored_key, stored_base = await _lookup_stored_byok_key(db, user_id, provider_hint=model_provider)
        if stored_key:
            effective_user_key = stored_key
            effective_base_url = stored_base

    if raw_model and raw_model.startswith("llamacpp/"):
        effective_user_key = None

    if effective_user_key:
        effective_base = effective_base_url or base_url
        client = AsyncOpenAI(api_key=effective_user_key, base_url=effective_base)
    elif base_url != _LLM_API_BASE or api_key != _LLM_API_KEY:
        client = AsyncOpenAI(api_key=api_key, base_url=base_url)
    else:
        client = _client

    # Circuit breaker protection
    from app.core.circuit_breaker import CircuitOpenError, get_circuit_breaker
    from app.core.metrics import record_llm_request

    provider_name = _get_provider_for_model(raw_model) or "deepseek"
    breaker = get_circuit_breaker(provider_name)
    llm_start = time.time()

    # Build messages and process all DB-dependent setup BEFORE saving user message.
    messages_for_llm = await _build_chat_messages(db, thread_id)
    # Append the new user message — it isn't in DB history yet
    messages_for_llm.append({"role": "user", "content": content})

    # ── Phase 5: resolve workspace_id for allowlist filtering ────────
    _thread_obj = await get_chat_thread(db, thread_id)
    _workspace_id = _thread_obj.workspace_id if _thread_obj else None

    # Phase 5: async tool list with workspace allowlist
    openai_tools = await _get_chat_openai_tools(db=db, workspace_id=_workspace_id)

    # ── Phase 2: pre-compute user scopes for tool authorization ──────
    _cached_user_scopes: set[str] | None = None
    _cached_user_role: str | None = None
    try:
        from app.models.user import User as UserModel

        _urow = (await db.execute(select(UserModel).where(UserModel.id == user_id))).scalar_one_or_none()
        if _urow is not None:
            _cached_user_scopes = set(getattr(_urow, "scopes", []) or [])
            _cached_user_role = getattr(_urow, "role", None)
    except Exception:
        logger.debug("send_message_to_llm: scope pre-fetch failed (non-fatal)", exc_info=True)

    if attachments:
        messages_for_llm = await _process_attachments(db, messages_for_llm, attachments, raw_model)

    if web_search:
        messages_for_llm = await _inject_web_search(messages_for_llm, content)

    # Save user message, commit, and immediately close the session to prevent
    # idle-in-transaction timeout during the LLM call
    await create_chat_message(db, thread_id, "user", content, user_id=user_id)
    await db.commit()
    await db.close()

    try:
        async with breaker.protect():
            # ── Tool-calling loop (non-streaming) ───────────────────
            total_prompt_tokens = 0
            total_completion_tokens = 0
            total_total_tokens = 0
            executed_tools: list[dict] = []  # P0-2: track tool calls for REST response

            for _round in range(_MAX_TOOL_ROUNDS):
                create_kwargs: dict = {
                    "model": model,
                    "messages": messages_for_llm,
                }
                if openai_tools:
                    create_kwargs["tools"] = openai_tools

                response = await client.chat.completions.create(**create_kwargs)
                assistant_message = response.choices[0].message

                # Accumulate token usage across all rounds
                if response.usage:
                    total_prompt_tokens += response.usage.prompt_tokens or 0
                    total_completion_tokens += response.usage.completion_tokens or 0
                    total_total_tokens += response.usage.total_tokens or 0

                # Check for tool calls
                if assistant_message.tool_calls:
                    # Add assistant message with tool_calls to history
                    messages_for_llm.append(
                        {
                            "role": "assistant",
                            "content": assistant_message.content,
                            "tool_calls": [
                                {
                                    "id": tc.id,
                                    "type": "function",
                                    "function": {
                                        "name": tc.function.name,
                                        "arguments": tc.function.arguments,
                                    },
                                }
                                for tc in assistant_message.tool_calls
                            ],
                        }
                    )

                    # Execute each tool and add results
                    for tc in assistant_message.tool_calls:
                        _tool_t0 = time.time()
                        tool_result = await _execute_tool_call(
                            tc.function.name,
                            tc.function.arguments,
                            user_id=user_id,
                            _user_scopes=_cached_user_scopes,
                            _user_role=_cached_user_role,
                        )
                        _record_tool_cost_fire_and_forget(
                            user_id=user_id,
                            tool_name=tc.function.name,
                            duration_ms=(time.time() - _tool_t0) * 1000.0,
                            workspace_id=_workspace_id,
                        )
                        # P0-2: track executed tool for REST response
                        executed_tools.append(
                            {
                                "tool": tc.function.name,
                                "call_id": tc.id,
                                "arguments": tc.function.arguments,
                                "result": tool_result,
                            }
                        )
                        messages_for_llm.append(
                            {
                                "role": "tool",
                                "tool_call_id": tc.id,
                                "content": tool_result,
                            }
                        )

                    continue  # Loop for next LLM call

                # No tool calls — final text response
                break
            else:
                # Loop exhausted — all rounds were tool calls
                assistant_message.content = (
                    assistant_message.content
                    or "I reached the maximum number of tool calls "
                    f"({_MAX_TOOL_ROUNDS}) without producing a final response."
                )

            # Fallback: retry without tools if the response is empty
            # (some models don't support function calling)
            if not (assistant_message.content or "").strip() and openai_tools:
                logger.info(
                    "send_message_to_llm: empty response with tools for %s, retrying without tools",
                    raw_model,
                )
                no_tools_response = await client.chat.completions.create(model=model, messages=messages_for_llm)
                if no_tools_response.choices:
                    assistant_message = no_tools_response.choices[0].message
                    response = no_tools_response

        assistant_content = assistant_message.content

        llm_duration = time.time() - llm_start
        record_llm_request(
            provider=provider_name,
            duration_seconds=llm_duration,
            prompt_tokens=total_prompt_tokens,
            completion_tokens=total_completion_tokens,
            success=True,
        )

        # Session was closed before LLM call to prevent idle-in-transaction timeout.
        # Always use fresh session for saving.
        try:
            await create_chat_message_fresh_session(thread_id, "assistant", assistant_content)
        except Exception as save_err:
            logger.warning("Assistant message save failed (non-fatal): %s", save_err)

        try:
            from app.services.usage_service import get_usage_service

            get_usage_service().record_usage(
                user_id=str(user_id),
                model_id=model,
                provider="byok" if user_api_key else "system",
                prompt_tokens=total_prompt_tokens,
                completion_tokens=total_completion_tokens,
                cost=total_total_tokens * 0.000002,
            )
        except Exception as ue:
            logger.warning("Usage recording failed (non-fatal): %s", ue)

        # ── P0-1: durable memory extraction via Celery ────────────────
        try:
            from app.tasks.memory_extraction_tasks import extract_memory_claims_task

            extract_memory_claims_task.delay(
                thread_id=thread_id,
                user_id=user_id,
                user_message=content,
                assistant_response=assistant_content,
            )
        except Exception:
            logger.debug("memory_extraction Celery dispatch failed (non-fatal)", exc_info=True)

        return {
            "success": True,
            "content": assistant_content,
            "tokens": total_total_tokens,
            "model": model,
            "tool_calls": executed_tools,
        }
    except CircuitOpenError as e:
        llm_duration = time.time() - llm_start
        record_llm_request(provider=provider_name, duration_seconds=llm_duration, success=False)
        logger.warning("Circuit breaker open for %s: %s", provider_name, e)
        return {
            "success": False,
            "content": f"Service temporarily unavailable ({provider_name}). Please try again later.",
            "tokens": 0,
            "model": model,
        }
    except Exception as e:
        llm_duration = time.time() - llm_start
        record_llm_request(provider=provider_name, duration_seconds=llm_duration, success=False)
        logger.error("send_message_to_llm failed: %s", e)
        return {"success": False, "content": str(e), "tokens": 0, "model": model}


async def _get_chat_openai_tools(
    db: AsyncSession | None = None,
    workspace_id: str | None = None,
) -> list[dict] | None:
    """Return chat-safe tools in OpenAI function-calling format, or None if none available.

    ADR-001: Computed Allowlist (Task 3.2 + P0-3)
    -----------------------------------------------
    The exposed tool set is computed as the intersection of 3 gates:

    1. **Visibility gate** (curation):  ``tool.metadata.visibility != "hidden"``
       Each tool declares visibility in its own ``ToolMetadata``:
       - ``default_on``: always exposed (Phase 1 core + sandboxd tools)
       - ``opt_in``:   exposed when available (Phase 2/3 read-only tools)
       - ``hidden``:   never exposed in chat (write ops, deferred tools)
    2. **Workspace gate** (existing):   workspace allowlist intersects
    3. **Scope gate** (existing, enforced in ``_execute_tool_call``):
       ``tool.metadata.required_scopes`` checked against user scopes

    Visibility tags are curation, NOT security.  ``required_scopes`` is
    the security boundary.  Adding a tool means tagging it in-file,
    not editing a central set.
    """
    try:
        from app.tools.base import get_tool_registry

        registry = get_tool_registry()

        # sandboxd tools gated by feature flag
        _SANDBOXD_IDS = frozenset(
            {
                "sandboxd_preview",
                "sandboxd_exec",
                "sandboxd_file_write",
                "sandboxd_file_read",
                "sandboxd_file_list",
                "sandboxd_serve",
                "browser_sandbox",
            }
        )

        # ── Compute exposed set ──────────────────────────────────────
        # P0-3: visibility now read from each tool's ToolMetadata.visibility
        # (set in-file). Default for untagged tools is "hidden" (the
        # ToolMetadata field default — see app/tools/base.py), so tools
        # not explicitly tagged are never exposed.
        exposed: list = []
        for tool in registry.list_all():
            vis = getattr(tool.metadata, "visibility", "hidden") or "hidden"
            if vis == "hidden":
                continue
            # Gate 1b: sandboxd tools require feature flag
            if tool.tool_id in _SANDBOXD_IDS and not settings.SANDBOXD_ENABLED:
                continue
            exposed.append(tool)

        # Gate 2: workspace allowlist intersection
        if db is not None and workspace_id:
            from app.models.workspace_models import get_workspace_tool_allowlist

            workspace_allowed = await get_workspace_tool_allowlist(db, workspace_id)
            if workspace_allowed is not None:
                exposed = [t for t in exposed if t.tool_id in workspace_allowed]

        # Gate 3 (required_scopes) is enforced in _execute_tool_call,
        # not here — we expose the tool schema so the LLM can call it,
        # and the scope check happens at execution time.

        tools = [t.to_openai_schema() for t in exposed]
        return tools or None
    except Exception:
        logger.debug("Failed to get chat tools from registry", exc_info=True)
        return None


async def _execute_tool_call(
    tool_name: str,
    arguments_json: str,
    user_id: int | None = None,
    workspace_id: str | None = None,
    _user_scopes: set[str] | None = None,
    _user_role: str | None = None,
) -> str:
    """Execute a single tool call via the registry and return the result as JSON.

    Phase 1: adds scope-based authorization check before ``tool.execute()``.
    If the tool declares ``required_scopes``, verify the caller holds them.
    Tools with no required_scopes are unrestricted.

    Uses a direct scope check (same pattern as v2/tools.py ``_user_has_scopes``)
    rather than issuing-then-verifying a fresh capability token, which would
    always pass and provide no real authorization.
    """
    try:
        from app.tools.base import get_tool_registry

        registry = get_tool_registry()
        tool = registry.get(tool_name)
        if tool is None:
            return json.dumps({"error": f"Unknown tool: {tool_name}"})

        # Phase 2: scope-based authorization check with cached user scopes.
        # If ``_user_scopes`` is provided (resolved once before the tool loop
        # in stream_message_to_llm), use it for proper scope resolution.
        # If not provided, fall back to admin-role bypass or blanket deny.
        if tool.metadata.required_scopes and user_id is not None:
            from app.core.auth_constants import ADMIN_ROLES

            # Admin/owner roles bypass scope checks
            if _user_role and _user_role in ADMIN_ROLES:
                pass  # full access
            elif _user_scopes is not None:
                missing = [s for s in tool.metadata.required_scopes if s not in _user_scopes]
                if missing:
                    logger.warning(
                        "Tool %s requires scopes %s — user missing %s",
                        tool_name,
                        tool.metadata.required_scopes,
                        missing,
                    )
                    return json.dumps(
                        {
                            "error": f"capability denied: tool '{tool_name}' requires scopes {tool.metadata.required_scopes} "
                            f"(missing: {missing})"
                        }
                    )
            else:
                # No cached scopes available — deny as defense-in-depth
                logger.warning(
                    "Tool %s requires scopes %s — denying (no cached scopes).",
                    tool_name,
                    tool.metadata.required_scopes,
                )
                return json.dumps(
                    {"error": f"capability denied: tool '{tool_name}' requires scopes {tool.metadata.required_scopes}"}
                )

        args = json.loads(arguments_json) if arguments_json else {}
        result = await tool.execute(args)
        if result.success:
            return json.dumps(result.result)
        return json.dumps({"error": result.error})
    except json.JSONDecodeError:
        return json.dumps({"error": f"Invalid JSON arguments: {arguments_json}"})
    except Exception as e:
        logger.exception("Tool execution failed: %s", tool_name)
        return json.dumps({"error": str(e)})


async def generate_thread_title(
    db: AsyncSession,
    thread_id: int,
) -> str | None:
    """Generate a 3-5 word title for a thread based on its first exchange.

    Uses the first user message and first assistant response to prompt a fast/cheap
    model for a concise title. Does NOT store any messages in the thread — this is
    a read-only operation that only updates the thread title.

    Returns the generated title or None if the thread has < 2 messages.
    """
    thread = await get_chat_thread(db, thread_id)
    if thread is None:
        return None

    # Fetch the first two messages (first user prompt + first assistant response)
    result = await db.execute(
        select(ChatMessage).where(ChatMessage.thread_id == thread_id).order_by(ChatMessage.id.asc()).limit(2)
    )
    first_messages = result.scalars().all()

    if len(first_messages) < 2:
        return None

    user_msg = first_messages[0]
    assistant_msg = first_messages[1]

    # Build a minimal prompt for title generation
    title_prompt = (
        "Generate a very short, descriptive title (3-5 words max) for a conversation "
        "that starts with the following exchange. Return ONLY the title, no quotes, "
        "no punctuation at the end, no explanation.\n\n"
        f"User: {user_msg.content[:300]}\n\n"
        f"Assistant: {assistant_msg.content[:300]}"
    )

    # Use the default client with the fastest/cheapest model available.
    # Set LLM_FAST_MODEL env var to override (e.g. "deepseek/deepseek-v4-flash").
    # Falls back to _LLM_MODEL (default: deepseek/deepseek-v4-flash).
    try:
        model_name = os.getenv("LLM_FAST_MODEL", _LLM_MODEL)
        # Strip provider prefix if present
        if "/" in model_name:
            model_name = model_name.split("/", 1)[1]

        response = await _client.chat.completions.create(
            model=model_name,
            messages=[{"role": "user", "content": title_prompt}],
            max_tokens=20,
            temperature=0.3,
        )

        title = response.choices[0].message.content
        if title:
            # Clean up: strip quotes, newlines, extra whitespace
            title = title.strip().strip("\"'").strip()
            # Take first line only
            title = title.split("\n")[0].strip()
            # Strip common model prefixes like "Title:", "Subject:", or markdown bold
            import re

            title = re.sub(r"^(Title|Subject|Topic):?\s*", "", title, flags=re.IGNORECASE)
            title = title.strip("*").strip()
            # Truncate to reasonable length
            if len(title) > 100:
                title = title[:97] + "..."

            if title:
                # Update the thread title
                thread.title = title
                await db.flush()
                logger.info("Auto-titled thread %d: %s", thread_id, title)
                return title

    except Exception as e:
        logger.warning("Auto-titling failed for thread %d: %s", thread_id, e)

    return None


async def create_chat_branch(
    db: AsyncSession,
    user_id: int,
    parent_thread_id: int,
    parent_message_id: int,
    title: str,
) -> ChatBranch:
    """Create a new branch: copies messages up to parent_message_id into a new thread."""

    # Verify parent thread exists and belongs to user
    parent_thread = await get_chat_thread(db, parent_thread_id)
    if parent_thread is None:
        raise ValueError("Parent thread not found")

    # Create new thread for the branch
    branch_thread = await create_chat_thread(db, user_id, parent_thread.username, title)

    # Copy messages up to and including parent_message_id
    all_msgs, _ = await get_chat_messages(db, parent_thread_id)
    msgs_to_copy = [m for m in all_msgs if m.id <= parent_message_id]
    for msg in msgs_to_copy:
        await create_chat_message(db, branch_thread.id, msg.role, msg.content, user_id=msg.user_id)

    # Create branch record
    branch = ChatBranch(
        thread_id=branch_thread.id,
        parent_thread_id=parent_thread_id,
        parent_message_id=parent_message_id,
        user_id=user_id,
        title=title,
    )
    db.add(branch)
    await db.flush()
    await db.refresh(branch)
    return branch


async def list_chat_branches(
    db: AsyncSession,
    parent_thread_id: int,
) -> list[ChatBranch]:
    """List all branches from a given thread."""
    result = await db.execute(select(ChatBranch).where(ChatBranch.parent_thread_id == parent_thread_id))
    return list(result.scalars().all())


async def get_chat_branch(db: AsyncSession, branch_id: int) -> ChatBranch | None:
    """Get a single branch by ID."""
    result = await db.execute(select(ChatBranch).where(ChatBranch.id == branch_id))
    return result.scalar_one_or_none()


async def delete_chat_branch(db: AsyncSession, branch_id: int) -> bool:
    """Delete a branch and its thread."""
    branch = await get_chat_branch(db, branch_id)
    if branch is None:
        return False
    # Delete the branch thread (cascades messages)
    await delete_chat_thread(db, branch.thread_id)
    # Delete the branch record
    await db.delete(branch)
    await db.flush()
    return True
